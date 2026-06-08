"""Claude 工具定义与实现（会邀 + 飞书空间维护 + PPT 生成 + Bag 分析）"""
from __future__ import annotations

import concurrent.futures
import json
import logging
import os
import platform
import re
import subprocess
import sys
import tempfile
import textwrap
import threading
import urllib.parse
from datetime import datetime, timedelta

import pytz
import requests as _requests

from feishu.client import FeishuClient
from feishu.mcp_client import FeishuProjectMCP

logger = logging.getLogger(__name__)
TZ_SHANGHAI = pytz.timezone("Asia/Shanghai")

# PPT 输出目录（可通过环境变量覆盖）
PPT_OUTPUT_DIR = os.environ.get("MOMENTA_PPT_DIR") or os.path.join(
    os.path.expanduser("~"), "Desktop", "Claude"
)

# PPT 模版路径（使用项目目录中的 Momenta 模版）
# tools.py 位于 personal-assistant/agent/，模版在上两级目录 AI Agent/
_AGENT_BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
PPT_TEMPLATE_PATH = os.environ.get(
    "PPT_TEMPLATE_PATH",
    os.path.join(_AGENT_BASE_DIR, "PPT模版.pptx"),
)

# Mviz Skill 目录
_SKILL_BASE_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))), ".agents", "skills"
)
_FEISHU_DOC_SCRIPTS_DIR = os.path.join(
    _AGENT_BASE_DIR, ".agents", "skills", "feishu-doc", "scripts"
)
_STREAM_CATALOG_PATH = os.path.join(
    _SKILL_BASE_DIR, "stream-helper", "data", "stream-catalog.json"
)

# Mviz topic echo API（各环境）
_TOPIC_ECHO_ENDPOINTS = {
    "EU": "https://instant-cpp.cla.eu.momenta.works/python/topic_echo",
    "JP": "https://instant-cpp.jp.momenta.works/cpp/editor/serialize",
    "ME": "https://instant-cpp.me.momenta.works/python/topic_echo",
    "CN": "https://instant-cpp.momenta.works/python/topic_echo",
}
_META_API_DOMAINS = {
    "EU": "mviz-api.cla.eu.momenta.works",
    "JP": "mviz-api.jp.momenta.works",
    "ME": "mviz-api.me.momenta.works",
    "CN": "mviz-api.momenta.works",
}
# 只需 md5 参数（不需要 storage_name/key）的存储桶
_ONLY_MD5_STORAGES = {
    "mdi-suzhou-obs", "default", "cla-cdi-dss-default",
    "refresh-small-data-sh-tos", "simulation-prod-copied-for-dss-recover",
    "swe-bench-ct-obs", "cla-dcs-ct-all-in-one-s3",
    "dcs-me-processed-data-test", "cla-dis-ingested-data-me-prod",
}

# ---- 内存缓存 ----
# stream-catalog.json 只读一次，避免每次 search/lookup 都做磁盘 I/O
_STREAM_CATALOG_CACHE: dict | None = None


def _load_stream_catalog() -> "dict | str":
    """读取 stream-catalog.json，模块级缓存，只读一次"""
    global _STREAM_CATALOG_CACHE
    if _STREAM_CATALOG_CACHE is not None:
        return _STREAM_CATALOG_CACHE
    if not os.path.exists(_STREAM_CATALOG_PATH):
        return (
            f"[错误] stream-catalog.json 未找到：{_STREAM_CATALOG_PATH}\n"
            "请确认 stream-helper skill 已安装。"
        )
    try:
        with open(_STREAM_CATALOG_PATH, encoding="utf-8") as f:
            _STREAM_CATALOG_CACHE = json.load(f)
        logger.info("stream-catalog.json 已加载，共 %d 个 stream",
                    len(_STREAM_CATALOG_CACHE.get("streams", [])))
        return _STREAM_CATALOG_CACHE
    except Exception as e:
        return f"[错误] 无法读取 stream-catalog.json: {e}"


# feishu-sync-cli 路径（pip install --user 时安装在用户目录下，不在 PATH）
def _find_feishu_sync_cli() -> list[str]:
    """返回调用 feishu-sync-cli 的命令（列表形式）"""
    candidates = [
        "feishu-sync-cli",
        os.path.join(os.path.expanduser("~"), "Library", "Python", "3.9", "bin", "feishu-sync-cli"),
        os.path.join(os.path.expanduser("~"), "Library", "Python", "3.10", "bin", "feishu-sync-cli"),
        os.path.join(os.path.expanduser("~"), "Library", "Python", "3.11", "bin", "feishu-sync-cli"),
        os.path.join(os.path.expanduser("~"), ".local", "bin", "feishu-sync-cli"),
    ]
    for c in candidates:
        try:
            subprocess.run([c, "--help"], capture_output=True, timeout=5)
            return [c]
        except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
            continue
    # 最终回退：用 python3 -m 调用
    return [sys.executable, "-m", "feishu_sync.skill"]

_FEISHU_SYNC_CMD: list[str] = _find_feishu_sync_cli()
# feishu-sync-cli 读写同一份 token 文件，多线程并发调用会产生竞态，用锁串行化
_FEISHU_CLI_LOCK = threading.Lock()

# 待追踪 RSVP 状态的日历事件（持久化至 JSON，重启不丢失）
# {event_id: {"calendar_id", "title", "owner_open_id", "attendees": {open_id: rsvp_status}, "created_at"}}
_PENDING_EVENTS_FILE = os.path.join(os.path.dirname(__file__), "..", ".pending_events.json")


def _load_pending_events() -> dict[str, dict]:
    try:
        with open(_PENDING_EVENTS_FILE, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _save_pending_events() -> None:
    try:
        with open(_PENDING_EVENTS_FILE, "w", encoding="utf-8") as f:
            json.dump(_PENDING_EVENTS, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.warning("保存 pending_events 失败: %s", e)


_PENDING_EVENTS: dict[str, dict] = _load_pending_events()

# ------------------------------------------------------------------ #
# 工具 JSON Schema 定义（提供给 Claude API）
# ------------------------------------------------------------------ #

TOOL_DEFINITIONS = [
    # ── 会邀模块 ──────────────────────────────────────────────────
    {
        "name": "search_users",
        "description": (
            "按姓名关键词搜索飞书用户。当用户提到需要邀请某人但没有 @mention 时使用。"
            "返回匹配的用户列表，每条包含 name、open_id。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "keyword": {
                    "type": "string",
                    "description": "要搜索的姓名关键词，如 '张三' 或 '小明'",
                }
            },
            "required": ["keyword"],
        },
    },
    {
        "name": "create_meeting",
        "description": (
            "创建飞书日历会议邀请，并向所有与会者发送邀请。"
            "会议创建成功后返回日历链接和会议信息。\n"
            "支持将部分与会者标记为 Optional（可选参加），在 optional_attendee_open_ids 中填写。\n"
            "外部用户（非本公司飞书账号）无法通过 search_users 查到 open_id，"
            "请将其邮箱地址填入 attendee_emails，系统会以第三方方式邀请。\n"
            "【重要】如果与会者来自跨公司飞书群（外部用户），"
            "请在调用本工具之前先一次性询问用户所有外部成员的邮箱，"
            "不要创建会议后再逐个询问。"
            "例如：\"我看到群里有几位外部成员，请一次提供他们的邮箱地址（多个可用逗号或换行分隔）\"。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "会议标题/主题"},
                "start_time": {
                    "type": "string",
                    "description": "会议开始时间，ISO 8601 格式，如 '2026-03-14T15:00:00'",
                },
                "end_time": {
                    "type": "string",
                    "description": "会议结束时间，ISO 8601 格式，如 '2026-03-14T16:00:00'",
                },
                "attendee_open_ids": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "必选与会者的飞书 open_id 列表（含会议发起人）。仅限内部飞书用户。",
                },
                "attendee_emails": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "外部用户（非本公司飞书账号）的邮箱列表，用于第三方日历邀请。内部用户请用 attendee_open_ids。",
                    "default": [],
                },
                "optional_attendee_open_ids": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "可选（Optional）与会者的飞书 open_id 列表，这些人在日历中标记为非必选",
                    "default": [],
                },
                "description": {
                    "type": "string",
                    "description": "会议议程或备注（可选）",
                    "default": "",
                },
                "location": {
                    "type": "string",
                    "description": "会议地点或视频链接（可选）",
                    "default": "",
                },
            },
            "required": ["title", "start_time", "end_time", "attendee_open_ids"],
        },
    },
    {
        "name": "query_availability",
        "description": (
            "查询多位飞书用户在指定日期的共同空闲时间段，用于选择合适的会议时间。"
            "返回所有人都空闲、且满足所需时长的时间窗口列表。"
            "建议在 create_meeting 前调用，先确认时间再发邀请。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "user_open_ids": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "要查询的所有与会者 open_id 列表（包含发起人自己）",
                },
                "date": {
                    "type": "string",
                    "description": "查询日期，格式 YYYY-MM-DD，默认今天",
                    "default": "",
                },
                "duration_minutes": {
                    "type": "integer",
                    "description": "所需会议时长（分钟），默认 60",
                    "default": 60,
                },
                "work_start": {
                    "type": "string",
                    "description": "工作时间开始，格式 HH:MM，默认 09:00",
                    "default": "09:00",
                },
                "work_end": {
                    "type": "string",
                    "description": "工作时间结束，格式 HH:MM，默认 18:00",
                    "default": "18:00",
                },
            },
            "required": ["user_open_ids"],
        },
    },

    {
        "name": "cancel_meeting",
        "description": (
            "取消（删除）已创建的飞书日历会议，并向与会者发送取消通知。\n"
            "方式一：提供 event_id + calendar_id（最精准，来自 create_meeting 返回结果）。\n"
            "方式二：提供 title 关键词（系统自动在近期会议中搜索匹配）。\n"
            "方式三：event_id 和 title 均不填，系统列出所有待处理会议供用户选择——"
            "**用户只说「取消会议」而未指明是哪个时，直接调用此工具（不填任何参数）即可。**\n"
            "当用户说「取消所有」或「全部取消」时，设置 cancel_all=true，直接取消所有匹配会议无需二次确认。\n"
            "当 tool 返回多个匹配并附带 event_id+calendar_id 时，再次调用本工具并传入对应值即可精确取消。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "event_id": {
                    "type": "string",
                    "description": "要取消的日历事件 ID（可选，来自 create_meeting 返回结果）",
                    "default": "",
                },
                "calendar_id": {
                    "type": "string",
                    "description": "事件所在的日历 ID（可选，有 event_id 时一起提供）",
                    "default": "",
                },
                "title": {
                    "type": "string",
                    "description": "会议标题关键词，用于自动搜索匹配（event_id 未知时使用）",
                    "default": "",
                },
                "reason": {
                    "type": "string",
                    "description": "取消原因，会通过 IM 通知与会者（可选）",
                    "default": "",
                },
                "cancel_all": {
                    "type": "boolean",
                    "description": "为 true 时，取消所有标题匹配的会议（不询问选择），适合用户明确说「取消所有」时使用",
                    "default": False,
                },
            },
            "required": [],
        },
    },

    # ── 飞书知识库模块（统一分发工具）──────────────────────────────
    {
        "name": "feishu_action",
        "description": (
            "操作飞书（Wiki + IM 消息 + 群聊记录 + PPTX 编辑）。通过 action 参数选择具体操作：\n"
            "• search — 按关键词搜索 Wiki，返回相关页面列表（标题+URL）\n"
            "• read_page — 读取指定飞书页面的完整内容（Markdown）\n"
            "• list_pages — 列出某个空间或父页面下的所有子页面\n"
            "• create_page — 在指定位置新建 Wiki 页面\n"
            "• edit_page — 替换页面中的指定文本片段（精确匹配）\n"
            "• move_page — 将页面移动到另一个目录下（需要 url + target_url）\n"
            "• send_message — 发送飞书 IM 文本消息：\n"
            "  - 发私信给用户：open_id + message\n"
            "  - 发群消息：chat_id + message（send_to_chat=true）\n"
            "  - @mention 他人：在 message 中写 <at user_id=\"对方open_id\">姓名</at>\n"
            "• recall_message — 撤回（删除）一条消息（需要 message_id）\n"
            "• list_groups — 获取群列表：含 Bot 当前在群、用户在群但 Bot 不在（已授权 OAuth 时）、Bot 历史群（已被移出）\n"
            "• create_group — 创建飞书群聊并拉入成员（需要 group_name + member_open_ids）。返回群 chat_id 和邀请链接\n"
            "• delete_group — 解散飞书群聊（需要 chat_id；Bot 必须是群主）\n"
            "• get_group_members — 获取指定群的全部成员列表（返回 open_id + 姓名），不依赖消息历史，用于组织全员会议时获取群成员\n"
            "• read_group_messages — 读取指定群或 P2P 私信的消息。"
            "支持两种时间模式：① hours=N 读最近 N 小时；② start_time+end_time 精确读取某时间段（格式 'YYYY-MM-DD HH:MM'）。"
            "Bot 不在该群时自动用用户 OAuth token 尝试读取（需已完成授权）。"
            "⚠️ 禁止用此工具来识别群成员，请改用 get_group_members\n"
            "• apply_mentions — 将 wiki 页面中的 @Name 纯文本替换为真正的飞书 @mention（需要 url + mention_map）\n"
            "• inspect_pptx — 列出 PPTX 每页所有形状的坐标（英寸）、文字预览，用于了解结构后再做精确修改\n"
            "• edit_pptx — 下载飞书云盘中的 PPTX，在原文件上做文本替换 + 形状位置/尺寸调整（保留格式），保存到本地。\n"
            "  修改形状位置时先用 inspect_pptx 获取 slide_index 和 shape_index，再传入 shape_updates。\n"
            "  需要 obj_token（从 read_page 结果末尾的【文件 token：xxx】获取）\n"
            "• create_task — 创建飞书任务并指定负责人/关注者。\n"
            "  内部用户（同公司飞书账号）：task_assignee_open_ids 填 open_id 列表（先用 search_users 查）。\n"
            "  外部用户：无法通过 open_id 分配，可将外部用户姓名写入 task_description，并用 send_message 单独通知。\n"
            "• create_doc — 将 Markdown 内容创建为飞书文档（docx），支持本地图片和 drawio 附件。\n"
            "  需要 title + content（Markdown 文本）；ref_url 可选（指定目标知识库父页面，不填则创建到默认位置）。\n"
            "• drawio_to_board — 将 drawio XML 转换为飞书画板（Board），可在飞书中可视化编辑。\n"
            "  需要 title + drawio_content（drawio XML 文本）；ref_url 可选（指定目标知识库父页面）。\n\n"
            "• group_summary — 按需触发群聊摘要（biweekly/monthly/quarterly）。"
            "日报/周报由系统自动触发，无需调用此工具。"
            "摘要将发布到飞书 Wiki 并发送 DM 通知。\n"
            "• calendar_events — 查询用户本周日历事件，返回会议列表（标题、时间、地点）。无需额外参数。\n\n"
            "典型工作流：search → read_page（了解内容）→ create_page / edit_page（写入）\n"
            "修改 PPTX 时序图工作流：read_page（获取 obj_token）→ inspect_pptx（查看形状坐标）→ edit_pptx（替换文字 + 移动形状）"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["search", "read_page", "list_pages", "create_page", "edit_page", "move_page", "send_message", "recall_message", "list_groups", "create_group", "delete_group", "get_group_members", "read_group_messages", "apply_mentions", "inspect_pptx", "edit_pptx", "create_task", "create_doc", "drawio_to_board", "group_summary", "calendar_events"],
                    "description": "操作类型",
                },
                "query": {
                    "type": "string",
                    "description": "搜索关键词，如 '日系项目方案'（action=search 时必填）",
                },
                "limit": {
                    "type": "integer",
                    "description": "搜索结果数量，默认 5（action=search 时可选）",
                    "default": 5,
                },
                "url": {
                    "type": "string",
                    "description": "页面 URL（action=read_page / list_pages / edit_page / move_page 时必填）",
                },
                "target_url": {
                    "type": "string",
                    "description": "目标父页面 URL（action=move_page 时必填，即移动到哪个目录下）",
                },
                "ref_url": {
                    "type": "string",
                    "description": "参考页面 URL（action=create_page 时必填）：position=child 时为父页面，position=sibling 时为同级页面",
                },
                "title": {
                    "type": "string",
                    "description": "新页面标题（action=create_page 时必填）",
                },
                "content": {
                    "type": "string",
                    "description": "页面内容（Markdown，action=create_page 时可选）",
                    "default": "",
                },
                "position": {
                    "type": "string",
                    "enum": ["child", "sibling"],
                    "description": "创建位置（action=create_page 时）：child=子页面（默认），sibling=同级页面",
                    "default": "child",
                },
                "old_string": {
                    "type": "string",
                    "description": "要替换的原始文本，需精确匹配（action=edit_page 时必填）",
                },
                "new_string": {
                    "type": "string",
                    "description": "替换后的新文本（action=edit_page 时必填）",
                },
                "open_id": {
                    "type": "string",
                    "description": "接收方的飞书 open_id（action=send_message 发私信时必填）",
                },
                "message": {
                    "type": "string",
                    "description": "要发送的文本消息内容（action=send_message 时必填）。可在文本中嵌入 <at user_id=\"open_id\">姓名</at> 来 @mention 他人。",
                },
                "send_to_chat": {
                    "type": "boolean",
                    "description": "是否发到群聊（action=send_message 时可选，true 时使用 chat_id 发群消息，false/不填时使用 open_id 发私信）",
                    "default": False,
                },
                "message_id": {
                    "type": "string",
                    "description": "消息 ID（action=recall_message 时必填）。格式如 om_xxxxxxxx。",
                },
                "chat_id": {
                    "type": "string",
                    "description": "群聊 ID（action=send_message 发群消息时必填；action=read_group_messages 时可选；可先用 list_groups 获取）",
                },
                "group_name": {
                    "type": "string",
                    "description": "新建群聊的名称（action=create_group 时必填，如 'XX项目讨论群'）",
                },
                "member_open_ids": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "要拉入群聊的成员 open_id 列表（action=create_group 时必填）。先用 search_users 获取 open_id",
                },
                "p2p_open_id": {
                    "type": "string",
                    "description": "P2P 私信对象的 open_id（action=read_group_messages 时可选，用于读取与某用户的私信历史；与 chat_id 二选一）",
                },
                "hours": {
                    "type": "integer",
                    "description": "读取最近多少小时的消息（action=read_group_messages 时可选，默认 24；当 start_time/end_time 均填写时本参数被忽略）",
                    "default": 24,
                },
                "start_time": {
                    "type": "string",
                    "description": (
                        "消息起始时间（action=read_group_messages 时可选）。"
                        "支持格式：'2025-03-20 14:00'、'2025-03-20T14:00:00+08:00'。"
                        "不含时区时默认按上海时间（UTC+8）解析。"
                        "与 end_time 同时填写时，忽略 hours 参数，精确读取该时间段消息。"
                    ),
                },
                "end_time": {
                    "type": "string",
                    "description": (
                        "消息截止时间（action=read_group_messages 时可选，与 start_time 配合使用）。"
                        "格式同 start_time。"
                    ),
                },
                "mention_map": {
                    "type": "object",
                    "description": "姓名 → open_id 映射，用于 apply_mentions（action=apply_mentions 时必填）。示例：{\"张三\": \"ou_xxx\", \"李四\": \"ou_yyy\"}",
                    "additionalProperties": {"type": "string"},
                },
                "obj_token": {
                    "type": "string",
                    "description": "飞书云盘文件 token（action=inspect_pptx / edit_pptx 时必填）。从 read_page 返回结果末尾的【文件 token：xxx】中获取",
                },
                "slide_index": {
                    "type": "integer",
                    "description": "幻灯片索引，0 起（action=inspect_pptx 时可选，默认返回全部页）",
                },
                "replacements": {
                    "type": "array",
                    "description": "文本替换列表（action=edit_pptx 时可选）。每项为 {\"old\": \"原文本\", \"new\": \"新文本\"}",
                    "items": {
                        "type": "object",
                        "properties": {
                            "old": {"type": "string"},
                            "new": {"type": "string"},
                        },
                        "required": ["old", "new"],
                    },
                },
                "shape_updates": {
                    "type": "array",
                    "description": (
                        "形状位置/尺寸更新列表（action=edit_pptx 时可选）。先用 inspect_pptx 获取 slide_index 和 shape_index，"
                        "再传入要修改的形状。left/top/width/height 均为英寸（浮点数），省略的字段保持不变。"
                        "示例：[{\"slide\": 2, \"shape\": 5, \"left\": 3.2}, {\"slide\": 2, \"shape\": 7, \"left\": 3.2, \"top\": 4.0}]"
                    ),
                    "items": {
                        "type": "object",
                        "properties": {
                            "slide": {"type": "integer", "description": "幻灯片索引（0 起）"},
                            "shape": {"type": "integer", "description": "形状索引（0 起，来自 inspect_pptx）"},
                            "left":   {"type": "number", "description": "距左边缘（英寸）"},
                            "top":    {"type": "number", "description": "距上边缘（英寸）"},
                            "width":  {"type": "number", "description": "宽度（英寸）"},
                            "height": {"type": "number", "description": "高度（英寸）"},
                        },
                        "required": ["slide", "shape"],
                    },
                },
                "output_filename": {
                    "type": "string",
                    "description": "保存的文件名（action=edit_pptx 时可选，默认为 '修改后.pptx'）",
                    "default": "修改后.pptx",
                },
                "task_title": {
                    "type": "string",
                    "description": "任务标题（action=create_task 时必填）",
                },
                "task_description": {
                    "type": "string",
                    "description": "任务详情/备注（action=create_task 时可选）。可在此注明无法通过 open_id 邀请的外部用户姓名",
                    "default": "",
                },
                "task_due": {
                    "type": "string",
                    "description": "截止时间（action=create_task 时可选），格式：'2026-03-21 18:00' 或 '2026-03-21'（上海时间）",
                },
                "task_assignee_open_ids": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "负责人的飞书 open_id 列表（action=create_task 时可选）。仅支持内部用户，先用 search_users 查询 open_id",
                    "default": [],
                },
                "task_follower_open_ids": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "关注者（抄送）的飞书 open_id 列表（action=create_task 时可选）。仅支持内部用户",
                    "default": [],
                },
                "drawio_content": {
                    "type": "string",
                    "description": "drawio XML 文本内容（action=drawio_to_board 时必填）",
                },
                "report_type": {
                    "type": "string",
                    "enum": ["biweekly", "monthly", "quarterly"],
                    "description": (
                        "摘要报告类型（action=group_summary 时必填）。"
                        "biweekly=双周报（≤30天有消息的群）、"
                        "monthly=月报（≤90天）、"
                        "quarterly=三月报（≤180天）。"
                        "daily/weekly 由系统自动触发，无需手动调用。"
                    ),
                },
                "days_back": {
                    "type": "integer",
                    "description": "回看天数（action=group_summary 时可选，默认由 report_type 决定）",
                },
            },
            "required": ["action"],
        },
    },

    # ── Bag / Mviz 分析模块 ─────────────────────────────────────────
    {
        "name": "bag_analysis",
        "description": (
            "分析 Mviz bag 数据，支持 stream 查询和 topic 数据下载。通过 action 参数选择操作：\n"
            "• search_streams — 按关键词搜索 Mviz stream（如'车辆速度'、'lidar'、'aes'）\n"
            "• lookup_stream — 查询指定 stream 的数据源（ROS topic）和渲染逻辑\n"
            "• resolve_url — 解析 Mviz URL，提取 bag MD5/storage 参数\n"
            "• download_topic — 从 Mviz URL/MD5 下载指定 ROS topic 数据并返回摘要\n\n"
            "典型工作流：search_streams 找目标 → lookup_stream 确认 topic → download_topic 获取数据 → 分析问题"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["search_streams", "lookup_stream", "resolve_url", "download_topic"],
                    "description": "操作类型",
                },
                "url": {
                    "type": "string",
                    "description": "Mviz URL 或 bag URL（resolve_url / download_topic 必填）",
                },
                "topic": {
                    "type": "string",
                    "description": "ROS topic 名称（download_topic 必填），如 /mla/egopose 或 /perception/fusion/object",
                },
                "stream_name": {
                    "type": "string",
                    "description": "Stream 精确名称（lookup_stream 必填），如 /fusion/object_detection/ego",
                },
                "query": {
                    "type": "string",
                    "description": "关键词（search_streams 必填），如 '车辆速度' 或 'aes_status'",
                },
                "bag_md5": {
                    "type": "string",
                    "description": "bag MD5（可选）。若已通过 resolve_url 获得，直接传入可跳过 URL 二次解析，节省时间",
                },
                "max_frames": {
                    "type": "integer",
                    "description": "下载时返回的最大帧数（默认 15，大 topic 建议设小一些）",
                    "default": 15,
                },
            },
            "required": ["action"],
        },
    },

    # ── 飞书项目 MCP 模块 ──────────────────────────────────────────
    {
        "name": "feishu_project",
        "description": (
            "通过 MCP 查询、管理和分析飞书项目工作项。通过 action 参数选择操作：\n"
            "• get_workitem_brief — 查询工作项基本信息（需要 project_key + work_item_id）\n"
            "• get_workitem_info — 查询工作项类型配置（字段列表、角色等，需要 project_key + work_item_type）\n"
            "• get_view_detail — 查询视图数据（需要 project_key + view_id，可选 fields、page_num）\n"
            "• create_workitem — 创建工作项（需要 project_key + work_item_type + fields_dict）\n"
            "• update_field — 更新工作项字段（需要 project_key + work_item_id + fields_updates）\n"
            "• finish_node — 完成工作项节点（需要 project_key + work_item_id + nodes）\n"
            "• get_node_detail — 查询节点详情（需要 project_key + work_item_id + node_id）\n"
            "• search_by_mql — 使用 MOQL 语句高级查询（需要 moql）\n"
            "• list_schedule — 查询用户日程（需要 project_key + user_keys + start_date + end_date）\n\n"
            "项目名自动映射：读取 skill 目录的 project-mapping.md 确定 project_key，无需手动查询。\n"
            "工作项类型别名：读取 workitem-type-mapping.md 映射（如'PPM需求'→'PPM导入'、'Bug'→'缺陷'）。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["get_workitem_brief", "get_workitem_info", "get_view_detail",
                             "create_workitem", "update_field", "finish_node",
                             "get_node_detail", "search_by_mql", "list_schedule"],
                    "description": "操作类型",
                },
                "project_key": {
                    "type": "string",
                    "description": "飞书项目的 project_key（大多数操作必填）。可通过读取 project-mapping.md 由项目名得到",
                },
                "work_item_id": {
                    "type": "string",
                    "description": "工作项 ID（get_workitem_brief / update_field / finish_node / get_node_detail 必填）",
                },
                "work_item_type": {
                    "type": "string",
                    "description": "工作项类型（get_workitem_info / create_workitem 必填，如 '需求'、'缺陷'、'任务'）",
                },
                "fields_list": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "要返回的字段名列表（get_workitem_brief / get_view_detail 可选，如 ['名称', '负责人', '状态']）",
                },
                "fields_dict": {
                    "type": "object",
                    "description": "创建工作项时的字段值字典（create_workitem 必填，如 {'名称': 'xxx', '负责人': 'user_key'}）",
                    "additionalProperties": True,
                },
                "fields_updates": {
                    "type": "array",
                    "description": "更新工作项字段的列表（update_field 必填），每项为 {field_key, field_value, field_type}",
                    "items": {"type": "object"},
                },
                "view_id": {
                    "type": "string",
                    "description": "视图 ID（get_view_detail 必填，从飞书项目视图 URL 中获取）",
                },
                "page_num": {
                    "type": "integer",
                    "description": "分页页码（get_view_detail 可选，默认 1）",
                    "default": 1,
                },
                "nodes": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "要完成的节点 ID 列表（finish_node 必填）",
                },
                "node_id": {
                    "type": "string",
                    "description": "节点 ID（get_node_detail 必填）",
                },
                "moql": {
                    "type": "string",
                    "description": "MOQL 查询语句（search_by_mql 必填），格式：SELECT ... FROM `project_key`.`work_item_type` WHERE ...",
                },
                "user_keys": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "用户 key 列表（list_schedule 必填）",
                },
                "start_date": {
                    "type": "string",
                    "description": "开始日期（list_schedule 必填，格式 YYYY-MM-DD）",
                },
                "end_date": {
                    "type": "string",
                    "description": "结束日期（list_schedule 必填，格式 YYYY-MM-DD）",
                },
            },
            "required": ["action"],
        },
    },

    # ── PPT 生成模块 ───────────────────────────────────────────────
    {
        "name": "generate_ppt",
        "description": (
            "根据内容大纲生成 Momenta 风格的客户汇报 PPT（.pptx 文件）。"
            "生成完成后返回文件路径，文件保存在桌面 Claude 文件夹下。"
            "调用前请先通过 feishu_action(search/read_page) 收集好相关资料，"
            "再整理为结构化的 slides 列表传入。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "PPT 封面主标题"},
                "subtitle": {
                    "type": "string",
                    "description": "PPT 封面副标题，如客户名称、日期等",
                    "default": "",
                },
                "slides": {
                    "type": "array",
                    "description": "幻灯片列表（封面之后的内容页）",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "本页标题"},
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "要点列表（每条一行）",
                            },
                            "note": {
                                "type": "string",
                                "description": "页面备注/说明（可选）",
                                "default": "",
                            },
                        },
                        "required": ["title", "bullets"],
                    },
                },
                "output_filename": {
                    "type": "string",
                    "description": "输出文件名（不含路径），默认为 '客户汇报.pptx'",
                    "default": "客户汇报.pptx",
                },
            },
            "required": ["title", "slides"],
        },
    },
    # ── PM 项目计划模块 ─────────────────────────────────────────────
    {
        "name": "generate_pm_plan",
        "description": (
            "生成 Momenta L2++ 适配项目二级开发计划甘特图（HTML + HD PNG）。"
            "包含标准里程碑（T0/T+4W/T+10W/T+18W/T+20W）和全套模块任务（底盘联调/控制匹配/"
            "相机/Radar/USS/Lidar/域控/CP主线对齐/AEB/APA/NOA/压测），可自定义覆盖。"
            "生成后可上传到飞书 Drive 文件夹，或在飞书 Wiki 创建项目计划文档。"
            "触发词：生成项目计划、排项目计划、生成二级计划、项目排期、PM计划。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "project_name": {
                    "type": "string",
                    "description": "项目名称，如「上汽 L2++ 适配项目」",
                },
                "t0_date": {
                    "type": "string",
                    "description": "T0 日期（EP 车辆到位），格式 YYYY-MM-DD",
                },
                "total_weeks": {
                    "type": "integer",
                    "description": "总周期（周数），默认 20",
                    "default": 20,
                },
                "drive_folder_token": {
                    "type": "string",
                    "description": "飞书 Drive 文件夹 token（可选），上传 HTML+PNG 到该文件夹",
                    "default": "",
                },
                "wiki_parent_url": {
                    "type": "string",
                    "description": "飞书 Wiki 父节点 URL（可选），在该节点下创建项目计划文档",
                    "default": "",
                },
                "milestones": {
                    "type": "array",
                    "description": (
                        "自定义里程碑列表（可选），不传则用标准 T0/T+4W/T+10W/T+18W/T+20W。"
                        "每项格式：{week: int, name: str, version: str}"
                    ),
                    "items": {
                        "type": "object",
                        "properties": {
                            "week": {"type": "integer"},
                            "name": {"type": "string"},
                            "version": {"type": "string"},
                        },
                        "required": ["week", "name"],
                    },
                },
                "modules_override": {
                    "type": "array",
                    "description": (
                        "覆盖默认模块列表（可选）。格式：[{name: str, tasks: [{label, start, end, color}]}]"
                    ),
                    "items": {"type": "object"},
                },
            },
            "required": ["project_name", "t0_date"],
        },
    },
    {
        "name": "search_gb_standard",
        "description": (
            "搜索国家标准（GB强标）原文条款，用于回答 L2++/L3/L4 自动驾驶相关国标问题。"
            "当用户询问以下内容时调用：国标/GB/强标/组合驾驶辅助/LACS/领航/L2++/L3/L4/"
            "ADS 的功能要求、技术参数、试验方法、条款内容、术语定义、强标符合性分析等。"
            "先用关键词搜索定位相关条款，再基于原文条款回答，不得凭记忆作答。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": (
                        "搜索关键词（支持中英文），如 '横向加速度'、'TTC'、'接管能力'、"
                        "'4.6.1'、'MRM'、'ODD'、'L3级' 等。建议用核心术语，避免长句。"
                    ),
                },
                "standard": {
                    "type": "string",
                    "enum": ["l2pp", "ads_l3l4"],
                    "description": (
                        "l2pp：《智能网联汽车 组合驾驶辅助系统安全要求》GB 报批稿（L2++/LACS/领航）；"
                        "ads_l3l4：《智能网联汽车 自动驾驶系统安全要求》征求意见稿（L3/L4 ADS）"
                    ),
                },
            },
            "required": ["query", "standard"],
        },
    },
    {
        "name": "debate",
        "description": (
            "让两个 AI 角色就某个话题进行多轮辩论/互相校验，生成结构化的讨论报告。\n"
            "适合场景：方案评审、决策校验、风险分析、技术方案对比、头脑风暴。\n"
            "流程：Agent A 分析 → Agent B 审查/质疑 → A 回应（可多轮）→ 仲裁者总结共识与结论。\n"
            "两个角色可以自定义，例如：'产品经理' vs '工程师'、'乐观派' vs '悲观派'、'支持方' vs '反对方'。\n"
            "结果可选择直接发送到飞书群（填 output_chat_id）。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "topic": {
                    "type": "string",
                    "description": "辩论/讨论的核心话题或问题，尽量具体",
                },
                "context": {
                    "type": "string",
                    "description": "背景信息（可选），如相关文档内容、数据、已知条件等",
                    "default": "",
                },
                "persona_a": {
                    "type": "string",
                    "description": "Agent A 的角色描述（可选，默认：资深分析师）。例如：'产品经理，关注用户价值和市场落地'",
                    "default": "",
                },
                "persona_b": {
                    "type": "string",
                    "description": "Agent B 的角色描述（可选，默认：批判性审查者）。例如：'资深工程师，关注技术可行性和实现成本'",
                    "default": "",
                },
                "rounds": {
                    "type": "integer",
                    "description": "辩论轮数（1-3，默认 2）。1轮适合快速校验，2-3轮适合深度讨论",
                    "default": 2,
                },
                "output_chat_id": {
                    "type": "string",
                    "description": "将完整讨论报告发到指定飞书群（可选，填 chat_id；可先用 feishu_action list_groups 获取）",
                    "default": "",
                },
            },
            "required": ["topic"],
        },
    },
    {
        "name": "find_skills",
        "description": (
            "搜索 Agent Skill 推荐。当用户想找某类功能的 skill、问「有没有能做 X 的工具/skill」时使用。"
            "同时搜索公共技能库（skills.sh）和 Momenta 内部技能库（mcp-hub.momenta.works），"
            "返回匹配的 skill 名称、描述和安装命令。"
            "触发词：找 skill、搜索 skill、有没有 X 的 skill/工具、推荐一个能做 X 的 skill。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": (
                        "搜索关键词，描述需要的功能。"
                        "公共库用英文关键词效果更好，如 'code review'、'diagram'、'deploy'；"
                        "内部库支持中文，如 '飞书'、'自动驾驶'、'测试用例'。"
                        "可以多次调用此工具，每次用不同关键词扩大搜索范围。"
                    ),
                },
                "source": {
                    "type": "string",
                    "enum": ["both", "public", "internal"],
                    "description": "搜索范围：both（默认）= 同时搜两个库；public = 仅公共库；internal = 仅内部库",
                    "default": "both",
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "np_aeb_calc",
        "description": (
            "NP&AEB 制动模型精确计算器。当用户提问涉及制动/刹车，且涉及距离计算（感知检出距离、制动距离、停车距离等）时使用。\n"
            "支持两种模式：\n"
            "  distance  ：已知初始车速(v0)、AEB介入速度(v_aeb)、NP制动边界 → 计算所需感知检出距离及各阶段分解\n"
            "  latest_aeb：已知感知检出距离(s_det)、初始车速(v0)、NP制动边界 → 反推最晚AEB介入速度\n"
            "NP limits 格式：[{\"v\": null, \"a\": 6, \"j\": 8}]（全速段），或按速度段分段：[{\"v\": \"50\", \"a\": 5, \"j\": 7}, {\"v\": \"100\", \"a\": 4, \"j\": 6}]，或阈值：[{\"v\": \">80\", \"a\": 3, \"j\": 5}]。\n"
            "默认参数：制动延迟 t_delay=0.5s，AEB减速度 aeb_a=9 m/s²，AEB jerk aeb_j=20 m/s³（无需用户提供时使用默认值）。\n"
            "注意：v0 为实际车速（km/h），非表显车速。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "mode": {
                    "type": "string",
                    "enum": ["distance", "latest_aeb"],
                    "description": "distance：正向计算所需检出距离；latest_aeb：反推最晚AEB介入速度",
                },
                "v0": {
                    "type": "number",
                    "description": "初始实际车速，单位 km/h",
                },
                "v_aeb": {
                    "type": "number",
                    "description": "AEB 介入速度，单位 km/h（distance 模式必填）",
                },
                "s_det": {
                    "type": "number",
                    "description": "感知检出距离，单位 m（latest_aeb 模式必填）",
                },
                "t_delay": {
                    "type": "number",
                    "description": "制动延迟，单位 s，默认 0.5",
                    "default": 0.5,
                },
                "aeb_a": {
                    "type": "number",
                    "description": "AEB 减速度幅值，单位 m/s²，默认 9",
                    "default": 9.0,
                },
                "aeb_j": {
                    "type": "number",
                    "description": "AEB jerk 幅值，单位 m/s³，默认 20",
                    "default": 20.0,
                },
                "limits": {
                    "type": "array",
                    "description": (
                        "NP 制动边界，JSON 数组。每项含 v（速度节点 km/h 或 null 表示全速段），a（最大减速度 m/s²），j（最大 jerk m/s³）。"
                        "例：[{\"v\": null, \"a\": 6, \"j\": 8}] 或 [{\"v\": \"50\", \"a\": 5, \"j\": 7}, {\"v\": \"100\", \"a\": 4, \"j\": 6}]"
                    ),
                    "items": {"type": "object"},
                },
            },
            "required": ["mode", "v0", "limits"],
        },
    },
    # ── 车辆 FMP 查询模块 ──────────────────────────────────────────
    {
        "name": "check_fmp_vehicles",
        "description": (
            "查询 FMP 平台上归属指定项目的车辆空闲状态。\n"
            "使用场景：用户问「Project-MXS现在有哪些空闲车」「哪辆车可以约」「FMP 上的车辆情况」等。\n"
            "返回空闲车辆列表和占用车辆列表。\n"
            "若 FMP session 未授权，会提示用户先运行登录脚本。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "project": {
                    "type": "string",
                    "description": "查询的项目/归属名称，如「Project-MXS」「Highway」，默认「Project-MXS」",
                    "default": "Project-MXS",
                },
                "hours": {
                    "type": "number",
                    "description": "查询未来几小时内的空闲情况，默认 8 小时",
                    "default": 8,
                },
            },
            "required": [],
        },
    },
    # ── 车辆预约模块 ────────────────────────────────────────────────
    {
        "name": "book_vehicle",
        "description": (
            "在「Project-MXS安全员协调」群里向 Fleet-Bot 发送车辆预约请求，完成预约并自动审批。\n"
            "使用场景：用户说「帮我预约 XXX 车」「约 Vehicle-C 那辆车」「预约车辆到22点」等车辆预约需求。\n"
            "发送后 Fleet-Bot 会自动处理预约流程，无需人工再点击确认。"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "vehicle_id": {
                    "type": "string",
                    "description": "车辆 ID，例如 Vehicle-C108620-A43902、Vehicle-P008620-B78738、Vehicle-C2-7836",
                },
                "time_range": {
                    "type": "string",
                    "description": "用车时间范围，例如「从现在到22:00」「今天14:00到22:00」「明天上午10点到晚上22点」",
                },
                "task_name": {
                    "type": "string",
                    "description": "用车任务名称，例如「集成测试」「数据采集」「执行器验收」，默认「集成测试」",
                    "default": "集成测试",
                },
                "project": {
                    "type": "string",
                    "description": "归属项目，例如「Project-MXS」。不填则不添加归属信息",
                    "default": "",
                },
            },
            "required": ["vehicle_id", "time_range"],
        },
    },
    # ── 故障码查询模块 ───────────────────────────────────────────────
    {
        "name": "query_bag_fault",
        "description": (
            "根据 VIN + 时间范围查询车辆 /mff_md/enable_signal_cmd 故障（enable_signal_cmd 信号）。\n"
            "从 ESS 查询触发事件，通过 pybagmining 按 topic 精准读取 bag 数据（无需下载整包），\n"
            "解析 error_code 并对照故障码映射表输出故障名称和受影响功能列表。\n"
            "典型场景：「查一下这辆车有没有故障」「这个 VIN 的 enable_signal_cmd 有没有 error」"
            "「帮我看看这车 HNP 为什么进不去」"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "vin": {
                    "type": "string",
                    "description": "17 位 VIN 码，例如 LFPH4ACL9T2A37871",
                },
                "time_points": {
                    "type": "string",
                    "description": (
                        "查询时间，多个时间点用空格分隔：\n"
                        "• YYYY-MM-DD — 整天\n"
                        "• YYYY-MM-DDTHH:MM:SS — 精确时刻（±5min窗口）\n"
                        "• 两个日期表示范围，如 '2026-05-15 2026-05-22'"
                    ),
                },
                "size": {
                    "type": "integer",
                    "description": "最多查询多少个 ESS 事件（默认 20，时间段较长时可调大）",
                    "default": 20,
                },
            },
            "required": ["vin", "time_points"],
        },
    },
]


# ------------------------------------------------------------------ #
# 工具执行
# ------------------------------------------------------------------ #

class ToolExecutor:
    def __init__(self, feishu: FeishuClient, message_cache=None, owner_open_id: str = "",
                 daily_summary_job=None):
        self.feishu = feishu
        self._message_cache = message_cache  # feishu.message_cache.MessageCache 实例（可为 None）
        self.owner_open_id = owner_open_id   # bot owner open_id，建群时自动加入
        self._daily_summary = daily_summary_job  # DailySummaryJob 实例，用于按需触发摘要
        # 同一会话内缓存 URL 解析结果，避免 resolve_url + download_topic 重复请求
        self._url_resolve_cache: dict[str, dict] = {}
        # 非 owner P2P 对话映射：sender_open_id → chat_id（由 main_ws 在收到消息时写入）
        self._p2p_chat_ids: dict[str, str] = {}

    def execute(self, tool_name: str, tool_input: dict) -> str:
        if tool_name == "search_users":
            return self._search_users(**tool_input)
        if tool_name == "query_availability":
            return self._query_availability(**tool_input)
        if tool_name == "create_meeting":
            return self._create_meeting(**tool_input)
        if tool_name == "cancel_meeting":
            return self._cancel_meeting(**tool_input)
        if tool_name == "feishu_action":
            return self._feishu_action(**tool_input)
        if tool_name == "bag_analysis":
            return self._bag_analysis(**tool_input)
        if tool_name == "generate_ppt":
            return self._generate_ppt(**tool_input)
        if tool_name == "generate_pm_plan":
            return self._generate_pm_plan(**tool_input)
        if tool_name == "feishu_project":
            return self._feishu_project(**tool_input)
        if tool_name == "search_gb_standard":
            return self._search_gb_standard(**tool_input)
        if tool_name == "debate":
            return self._debate(**tool_input)
        if tool_name == "find_skills":
            return self._find_skills(**tool_input)
        if tool_name == "np_aeb_calc":
            return self._np_aeb_calc(**tool_input)
        if tool_name == "book_vehicle":
            return self._book_vehicle(**tool_input)
        if tool_name == "check_fmp_vehicles":
            return self._check_fmp_vehicles(**tool_input)
        if tool_name == "query_bag_fault":
            return self._query_bag_fault(**tool_input)
        return f"未知工具: {tool_name}"

    # ── 会邀 ──────────────────────────────────────────────────────

    def _search_users(self, keyword: str) -> str:
        users = self.feishu.search_users(keyword, page_size=5)
        if not users:
            return f"未找到匹配 '{keyword}' 的用户。请确认姓名或让发起人直接 @mention。"
        lines = [f"找到 {len(users)} 个匹配用户："]
        for u in users:
            dept = f"（{u['department']}）" if u.get("department") else ""
            lines.append(f"- {u['name']}{dept}，open_id: {u['open_id']}")
        return "\n".join(lines)

    def _query_availability(
        self,
        user_open_ids: list[str],
        date: str = "",
        duration_minutes: int = 60,
        work_start: str = "09:00",
        work_end: str = "18:00",
    ) -> str:
        """查询共同空闲时间。使用 feishu-sync 获取个人日历忙碌时段。"""
        # 解析目标日期
        if not date:
            target_date = datetime.now(TZ_SHANGHAI).date()
        else:
            try:
                target_date = datetime.strptime(date, "%Y-%m-%d").date()
            except ValueError:
                return f"日期格式错误: {date}，请使用 YYYY-MM-DD"

        wh, wm = map(int, work_start.split(":"))
        eh, em = map(int, work_end.split(":"))
        day_start = TZ_SHANGHAI.localize(
            datetime(target_date.year, target_date.month, target_date.day, wh, wm)
        )
        day_end = TZ_SHANGHAI.localize(
            datetime(target_date.year, target_date.month, target_date.day, eh, em)
        )

        # 使用 feishu-sync 获取个人日历事件（替代已失效的 freebusy API）
        freebusy: dict[str, list[dict]] = {}
        try:
            raw = self._run_feishu_cli("calendar_events", "--json", timeout=20)
            if not raw.startswith("[飞书"):
                events = json.loads(raw)
                busy_slots = []
                for ev in events:
                    try:
                        ts_start = int(ev["start_time"]["timestamp"])
                        ts_end = int(ev["end_time"]["timestamp"])
                        ev_start = datetime.fromtimestamp(ts_start, tz=TZ_SHANGHAI)
                        ev_end = datetime.fromtimestamp(ts_end, tz=TZ_SHANGHAI)
                        if ev_start.date() == target_date and ev.get("free_busy_status") == "busy":
                            busy_slots.append({"start": ev_start.isoformat(), "end": ev_end.isoformat()})
                    except Exception:
                        pass
                # 用第一个 open_id 代表本人（拥有个人日历访问权限）
                if user_open_ids:
                    freebusy[user_open_ids[0]] = busy_slots
        except Exception as e:
            logger.warning("feishu-sync calendar_events 失败: %s", e)

        other_users = user_open_ids[1:] if len(user_open_ids) > 1 else []

        # 合并所有用户的忙碌区间
        all_busy: list[tuple[datetime, datetime]] = []
        for busy_list in freebusy.values():
            for b in busy_list:
                try:
                    bs = datetime.fromisoformat(b["start"])
                    be = datetime.fromisoformat(b["end"])
                    if bs.tzinfo is None:
                        bs = TZ_SHANGHAI.localize(bs)
                    if be.tzinfo is None:
                        be = TZ_SHANGHAI.localize(be)
                    all_busy.append((bs, be))
                except Exception:
                    pass

        all_busy.sort(key=lambda x: x[0])
        merged: list[list[datetime]] = []
        for bs, be in all_busy:
            if merged and bs <= merged[-1][1]:
                merged[-1][1] = max(merged[-1][1], be)
            else:
                merged.append([bs, be])

        # 在工作时间内找出空闲窗口
        duration = timedelta(minutes=duration_minutes)
        suggestions: list[tuple[datetime, datetime]] = []
        cursor = day_start
        for busy_s, busy_e in merged:
            if cursor + duration <= busy_s:
                # 当前空隙足够，每 30 分钟提一个候选时间
                slot = cursor
                while slot + duration <= busy_s and len(suggestions) < 6:
                    suggestions.append((slot, slot + duration))
                    slot += timedelta(minutes=30)
            cursor = max(cursor, busy_e)
        # 最后一段空隙
        slot = cursor
        while slot + duration <= day_end and len(suggestions) < 6:
            suggestions.append((slot, slot + duration))
            slot += timedelta(minutes=30)

        note = ""
        if other_users:
            note = f"\n（注：仅基于你的日历，{len(other_users)} 位邀请者的日历无法自动查询，请自行确认他们的空闲。）"

        if not suggestions:
            return (
                f"{target_date} 全天在工作时间 {work_start}-{work_end} 内，"
                f"找不到 {duration_minutes} 分钟的空闲时间段。{note}"
            )

        lines = [
            f"{target_date} 你的空闲时段（会议时长 {duration_minutes} 分钟）：",
        ]
        for i, (s, e) in enumerate(suggestions, 1):
            lines.append(f"  {i}. {s.strftime('%H:%M')} - {e.strftime('%H:%M')}")
        if note:
            lines.append(note)
        lines.append("\n请确认时间后，告诉我开始时间，我来创建会议并发送邀请。")
        return "\n".join(lines)

    def _create_meeting(
        self,
        title: str,
        start_time: str,
        end_time: str,
        attendee_open_ids: list[str],
        description: str = "",
        location: str = "",
        optional_attendee_open_ids: list[str] | None = None,
        attendee_emails: list[str] | None = None,
    ) -> str:
        try:
            start_dt = _parse_dt(start_time)
            end_dt = _parse_dt(end_time)
        except ValueError as e:
            return f"时间格式解析失败: {e}，请使用 ISO 8601 格式（如 2026-03-14T15:00:00）"

        optional_ids = optional_attendee_open_ids or []
        ext_emails = [e.strip() for e in (attendee_emails or []) if e.strip()]

        # Step 1: 创建事件（不含与会者）
        result = self.feishu.create_calendar_event(
            title=title,
            description=description,
            start_dt=start_dt,
            end_dt=end_dt,
            location=location,
        )
        if not result:
            return "会议创建失败，请检查飞书应用权限（需要 calendar:calendar 权限）。"

        event_id = result["event_id"]
        calendar_id = result["calendar_id"]

        # Step 2: 通过 /attendees API 发送正式日历邀请（含 optional 标记）
        all_ids = list(attendee_open_ids) + [oid for oid in optional_ids if oid not in attendee_open_ids]
        invite_ok, failed_oids = self.feishu.add_event_attendees(
            calendar_id, event_id, attendee_open_ids,
            optional_open_ids=optional_ids,
            attendee_emails=ext_emails,
        )

        # Step 3: IM 消息兜底通知（仅内部飞书用户；外部用户收日历邮件即可）
        start_str = start_dt.strftime("%Y-%m-%d %H:%M")
        end_str = end_dt.strftime("%H:%M")
        optional_set = set(optional_ids)
        for oid in all_ids:
            role_hint = "（Optional，可选参加）" if oid in optional_set else ""
            backup_msg = (
                f"【会议邀请{role_hint}】{title}\n"
                f"时间：{start_str} - {end_str}\n"
                f"地点：{location or '未设置'}\n"
                "请在飞书日历「邀请」栏查收并确认。"
            )
            try:
                self.feishu.send_text_to_user(oid, backup_msg)
            except Exception as e:
                logger.warning("发送备份消息给 %s 失败: %s", oid, e)

        # Step 4: 登记事件，供 RSVP 轮询任务追踪
        owner_open_id = getattr(self, "owner_open_id", "")
        _PENDING_EVENTS[event_id] = {
            "calendar_id": calendar_id,
            "title": title,
            "owner_open_id": owner_open_id,
            "attendees": {oid: "needs_action" for oid in all_ids},
            "created_at": datetime.now(TZ_SHANGHAI).timestamp(),
        }
        _save_pending_events()

        invite_note = (
            "日历邀请已发送（请在飞书日历「邀请」栏确认）"
            if invite_ok
            else "⚠️ 日历邀请发送失败（可能需审批 calendar:acl:write 权限），已发送飞书消息通知"
        )
        optional_note = f"\n可选参与者（Optional）：{len(optional_ids)} 人" if optional_ids else ""
        ext_note = f"\n外部用户（邮件邀请）：{', '.join(ext_emails)}" if ext_emails else ""
        im_note = "已向内部受邀者发送飞书消息通知，将实时追踪回复状态。" if all_ids else ""
        # 外部/跨租户用户无法通过 open_id 邀请，提示 Claude 询问邮箱
        failed_note = (
            f"\n⚠️ {len(failed_oids)} 名与会者被日历系统静默拒绝（跨公司/外部飞书用户无法通过 open_id 邀请）。\n"
            f"请**在一条消息中**列出这 {len(failed_oids)} 人的姓名，一次性请用户提供他们所有人的邮箱地址，"
            f"不要逐一询问（例如：\"以下参会者是外部用户，请一次提供所有邮箱：[姓名1]、[姓名2]...\"）。\n"
            f"收到全部邮箱后，用 attendee_emails=[...邮箱列表...] 重新调用 create_meeting，"
            f"同时在 attendee_open_ids 中仅保留本次成功添加的内部用户。\n"
            f"未添加的 open_ids: {failed_oids}"
        ) if failed_oids else ""
        return (
            f"会议已创建成功！\n"
            f"标题：{title}\n"
            f"时间：{start_str} - {end_str}\n"
            f"地点：{location or '未设置'}\n"
            f"必选与会人数：{len(attendee_open_ids)} 人{optional_note}{ext_note}\n"
            f"{invite_note}\n"
            f"{im_note}"
            f"{failed_note}\n"
            f"[event_id={event_id} calendar_id={calendar_id}]"
        )

    def _cancel_meeting(
        self,
        event_id: str = "",
        calendar_id: str = "",
        title: str = "",
        reason: str = "",
        cancel_all: bool = False,
    ) -> str:
        """取消（删除）日历事件，并通过 IM 通知与会者。

        优先用 event_id 精确匹配；若无 event_id，按 title 关键词搜索近期会议。
        """
        # ── 情况 1：有 event_id，直接操作 ──────────────────────────
        if event_id and calendar_id:
            return self._do_cancel(event_id, calendar_id, reason)

        # ── 情况 2：只有 title，先从持久化记录里精确查 ──────────────
        if title:
            kw = title.lower()
            matched = [
                (eid, info) for eid, info in _PENDING_EVENTS.items()
                if kw in info.get("title", "").lower()
            ]
            if len(matched) == 1:
                eid, info = matched[0]
                return self._do_cancel(eid, info["calendar_id"], reason, known_info=info)
            if len(matched) > 1:
                if cancel_all:
                    results = []
                    for eid, info in matched:
                        results.append(self._do_cancel(eid, info["calendar_id"], reason, known_info=info))
                    return "\n\n".join(results)
                titles = "\n".join(
                    f"- {info['title']} (event_id={eid} calendar_id={info['calendar_id']})"
                    for eid, info in matched
                )
                return (
                    f"找到 {len(matched)} 个匹配的会议，请告诉我取消哪个（或说「全部取消」）：\n{titles}"
                )

            # 持久化记录里没有，查询日历 API（近 30 天内）
            now = datetime.now(TZ_SHANGHAI)
            events = self.feishu.list_calendar_events(
                time_min=now - timedelta(days=1),
                time_max=now + timedelta(days=30),
                title_keyword=title,
            )
            if not events:
                return (
                    f"在近期日历中未找到标题含「{title}」的会议。\n"
                    "请提供 event_id 或确认会议标题是否正确。"
                )
            if len(events) == 1:
                e = events[0]
                return self._do_cancel(e["event_id"], e["calendar_id"], reason)
            # 多个匹配
            if cancel_all:
                results = []
                for e in events:
                    results.append(self._do_cancel(e["event_id"], e["calendar_id"], reason))
                return "\n\n".join(results)
            lines = "\n".join(
                f"- {e['title']} (event_id={e['event_id']} calendar_id={e['calendar_id']})"
                for e in events
            )
            return (
                f"找到 {len(events)} 个匹配的会议，请告诉我取消哪个（或说「全部取消」）：\n{lines}"
            )

        # ── 情况 3：无任何条件，列出所有待处理会议供用户选择 ──────────
        if _PENDING_EVENTS:
            lines = "\n".join(
                f"- {info.get('title', '(无标题)')}（event_id={eid}）"
                for eid, info in _PENDING_EVENTS.items()
            )
            return f"当前待处理会议如下，请告诉我取消哪个（或说「全部取消」）：\n{lines}"
        # 持久化记录为空，查询日历 API（近 30 天）
        now = datetime.now(TZ_SHANGHAI)
        events = self.feishu.list_calendar_events(
            time_min=now - timedelta(days=1),
            time_max=now + timedelta(days=30),
        )
        if not events:
            return "近期日历中未找到你创建的会议。"
        lines = "\n".join(
            f"- {e['title']}（{e.get('start', '')}，event_id={e['event_id']}）"
            for e in events
        )
        return f"找到以下近期会议，请告诉我取消哪个（或说「全部取消」）：\n{lines}"

    def _do_cancel(
        self,
        event_id: str,
        calendar_id: str,
        reason: str,
        known_info: dict | None = None,
    ) -> str:
        """执行实际删除并通知与会者"""
        event_info = known_info or _PENDING_EVENTS.get(event_id, {})
        title_text = event_info.get("title", "该会议")
        attendee_oids = list(event_info.get("attendees", {}).keys())

        ok = self.feishu.delete_calendar_event(calendar_id, event_id)
        if not ok:
            return f"取消会议失败（event_id={event_id}），请确认 event_id 正确或手动在飞书日历中删除。"

        _PENDING_EVENTS.pop(event_id, None)
        _save_pending_events()

        cancel_msg = (
            f"【会议取消】{title_text}\n该会议已被取消。"
            + (f"\n原因：{reason}" if reason else "")
        )
        for oid in attendee_oids:
            try:
                self.feishu.send_text_to_user(oid, cancel_msg)
            except Exception as e:
                logger.warning("发送取消通知给 %s 失败: %s", oid, e)

        return (
            f"会议「{title_text}」已成功取消。\n"
            f"已向 {len(attendee_oids)} 位与会者发送取消通知。"
            + (f"\n取消原因：{reason}" if reason else "")
        )

    # ── 飞书知识库 ────────────────────────────────────────────────

    def _feishu_action(
        self,
        action: str,
        query: str = "",
        limit: int = 5,
        url: str = "",
        target_url: str = "",
        ref_url: str = "",
        title: str = "",
        content: str = "",
        position: str = "child",
        old_string: str = "",
        new_string: str = "",
        open_id: str = "",
        message: str = "",
        send_to_chat: bool = False,
        message_id: str = "",
        chat_id: str = "",
        p2p_open_id: str = "",
        hours: int = 24,
        start_time: str = "",
        end_time: str = "",
        mention_map: dict | None = None,
        obj_token: str = "",
        slide_index: int | None = None,
        replacements: list | None = None,
        shape_updates: list | None = None,
        output_filename: str = "修改后.pptx",
        group_name: str = "",
        member_open_ids: list | None = None,
        task_title: str = "",
        task_description: str = "",
        task_due: str = "",
        task_assignee_open_ids: list | None = None,
        task_follower_open_ids: list | None = None,
        drawio_content: str = "",
        report_type: str = "",
        days_back: int | None = None,
    ) -> str:
        """统一分发飞书知识库操作"""
        if action == "search":
            if not query:
                return "错误：search 操作需要提供 query 参数。"
            return self._search_feishu(query, limit)
        if action == "read_page":
            if not url:
                return "错误：read_page 操作需要提供 url 参数。"
            return self._read_feishu_page(url)
        if action == "list_pages":
            if not url:
                return "错误：list_pages 操作需要提供 url 参数。"
            return self._list_feishu_pages(url)
        if action == "create_page":
            if not ref_url or not title:
                return "错误：create_page 操作需要提供 ref_url 和 title 参数。"
            return self._create_feishu_page(ref_url, title, content, position)
        if action == "edit_page":
            if not url or not old_string or not new_string:
                return "错误：edit_page 操作需要提供 url、old_string、new_string 参数。"
            return self._edit_feishu_page(url, old_string, new_string)
        if action == "move_page":
            if not url or not target_url:
                return "错误：move_page 操作需要提供 url（要移动的页面）和 target_url（目标父目录）参数。"
            return self._move_feishu_page(url, target_url)
        if action == "send_message":
            if not message:
                return "错误：send_message 操作需要提供 message 参数。"
            # 发群消息
            if send_to_chat or chat_id:
                if not chat_id:
                    return "错误：发群消息需要提供 chat_id 参数。"
                try:
                    mid = self.feishu.send_text_to_chat(chat_id, message)
                    if mid:
                        return f"群消息已发送。message_id={mid}（如需撤回可用此 ID 调用 recall_message）"
                    return "群消息发送失败，请检查 chat_id 是否正确。"
                except _requests.HTTPError as e:
                    return f"群消息发送异常（HTTP {e.response.status_code if e.response else '?'}）：{e}"
                except Exception as e:
                    return f"群消息发送异常：{e}"
            # 发私信给用户
            if not open_id:
                return "错误：发私信需要提供 open_id 参数；发群消息需要提供 chat_id 参数。"
            # 防止向 Bot 自身发消息（会导致 400）
            if not hasattr(self, "_bot_open_id_cache"):
                self._bot_open_id_cache = self.feishu.get_bot_open_id()
            if open_id == self._bot_open_id_cache:
                return "错误：接收人是 Bot 自身，无法向自己发消息。请确认接收人 open_id 是否正确。"
            try:
                mid = self.feishu.send_text_to_user(open_id, message)
                if mid:
                    return f"消息已发送。message_id={mid}（如需撤回可用此 ID 调用 recall_message）"
                return "消息发送失败，请检查 open_id 是否正确。"
            except _requests.HTTPError as e:
                if e.response is not None and e.response.status_code == 400:
                    return (
                        "消息发送失败：对方尚未与 Bot 建立私聊。"
                        "需要对方先在飞书中打开 Bot 的聊天窗口发送一条消息，之后 Bot 才能主动给对方发消息。"
                    )
                return f"消息发送异常（HTTP {e.response.status_code if e.response else '?'}）：{e}"
            except Exception as e:
                return f"消息发送异常：{e}"
        if action == "recall_message":
            if not message_id:
                return "错误：recall_message 操作需要提供 message_id 参数（格式如 om_xxxxxxxx）。"
            try:
                ok = self.feishu.recall_message(message_id)
                return "消息已撤回。" if ok else "消息撤回失败，请检查 message_id 是否正确或消息是否可撤回。"
            except _requests.HTTPError as e:
                status = e.response.status_code if e.response else "?"
                if status == 403:
                    return "消息撤回失败：无权限撤回此消息（只能撤回 Bot 自己发送的消息）。"
                return f"消息撤回异常（HTTP {status}）：{e}"
            except Exception as e:
                return f"消息撤回异常：{e}"
        if action == "list_groups":
            try:
                # Bot 当前已加入的群
                current_groups = self.feishu.get_joined_groups()
                current_ids = {g["chat_id"] for g in current_groups}
                lines = []
                for g in current_groups:
                    lines.append(f"- [Bot在群] {g.get('name', '未知群名')}（chat_id: {g['chat_id']}）")

                # 用户身份加入的群（Bot 不在，但用户有 OAuth 授权）
                user_groups = self.feishu.get_user_joined_groups()
                user_only_count = 0
                for g in user_groups:
                    cid = g["chat_id"]
                    if cid in current_ids:
                        continue
                    current_ids.add(cid)
                    lines.append(f"- [用户在群/Bot不在] {g.get('name', '未知群名')}（chat_id: {cid}）")
                    user_only_count += 1

                # 从本地消息缓存中补充：Bot 曾参与但已离开的群
                if self._message_cache:
                    for c in self._message_cache.get_known_chats():
                        cid = c["chat_id"]
                        if cid in current_ids:
                            continue
                        current_ids.add(cid)
                        label = "[历史群-已离开]" if c.get("chat_type") != "p2p" else "[P2P]"
                        name = c.get("name") or cid
                        lines.append(f"- {label} {name}（chat_id: {cid}）")

                # 从持久化 P2P 映射中补充私信联系人
                for sender_oid, cid in self._p2p_chat_ids.items():
                    if cid in current_ids or any(cid in l for l in lines):
                        continue
                    name = self.feishu._resolve_sender_name(sender_oid)
                    lines.append(f"- [P2P] 与 {name} 的私信（chat_id: {cid}，联系人 open_id: {sender_oid}）")

                if not lines:
                    return "Bot 当前不在任何群组中，也没有历史消息缓存。"
                summary = f"Bot 参与的会话（Bot 在群：{len(current_groups)} 个"
                if user_only_count:
                    summary += f"，用户在群/Bot不在：{user_only_count} 个"
                summary += "）："
                return summary + "\n" + "\n".join(lines)
            except Exception as e:
                return f"获取群列表失败：{e}"
        if action == "create_group":
            if not group_name:
                return "错误：create_group 需要提供 group_name（群名称）参数。"
            members = list(member_open_ids or [])
            # 确保 bot owner 始终在群里（即使 Claude 没有显式传入）
            if self.owner_open_id and self.owner_open_id not in members:
                members.append(self.owner_open_id)
            if not members:
                return "错误：create_group 需要提供 member_open_ids（成员 open_id 列表），至少包含 1 人。"
            try:
                result = self.feishu.create_group_chat(group_name, members)
                if not result.get("chat_id"):
                    err = result.get("error", "未知错误")
                    return f"建群失败：{err}（群名：{group_name}，成员数：{len(members)}）"
                chat_id_created = result["chat_id"]
                invite_link = result.get("invite_link", "")
                lines = [f"群「{group_name}」创建成功！", f"chat_id：{chat_id_created}"]
                if invite_link:
                    lines.append(f"邀请链接：{invite_link}")
                owner_note = "（含你自己）" if self.owner_open_id in members else ""
                lines.append(f"已拉入 {len(members)} 位成员{owner_note}。")

                # 向 bot owner 发送 DM，确保群入口可达（Bot 创建的群不会直接出现在成员列表里）
                if self.owner_open_id:
                    try:
                        if invite_link:
                            dm_msg = f"【新群已建】「{group_name}」\n点击加入：{invite_link}"
                        else:
                            dm_msg = f"【新群已建】「{group_name}」\nchat_id：{chat_id_created}\n（在飞书搜索群 ID 或群名即可找到）"
                        self.feishu.send_text_to_user(self.owner_open_id, dm_msg)
                        lines.append("✓ 已发送入群通知到你的私信")
                    except Exception as e:
                        logger.warning("发送建群 DM 失败: %s", e)
                return "\n".join(lines)
            except Exception as e:
                return f"建群失败：{e}"
        if action == "delete_group":
            if not chat_id:
                return "错误：delete_group 需要提供 chat_id 参数。"
            try:
                result = self.feishu.delete_group_chat(chat_id)
                if result.get("ok"):
                    return f"群（chat_id：{chat_id}）已成功解散。"
                return f"解散群失败：{result.get('error', '未知错误')}（Bot 需要是群主才能解散群）"
            except Exception as e:
                return f"解散群失败：{e}"
        if action == "get_group_members":
            if not chat_id:
                return "错误：get_group_members 需要提供 chat_id 参数。"
            try:
                members = self.feishu.get_chat_members(chat_id)
                if not members:
                    return "该群没有成员，或 Bot 无权限读取群成员列表（需要 im:chat:readonly 权限）。"
                lines = [f"- {m['name']}（open_id: {m['open_id']}）" for m in members]
                return f"群内成员共 {len(members)} 人：\n" + "\n".join(lines)
            except Exception as e:
                return f"获取群成员失败：{e}"
        if action == "read_group_messages":
            # 支持 p2p_open_id 作为 chat_id 的替代（从缓存中查找 P2P chat_id）
            if not chat_id and p2p_open_id:
                chat_id = self._p2p_chat_ids.get(p2p_open_id, "")
                if not chat_id:
                    return (
                        f"未找到与用户 {p2p_open_id} 的 P2P 对话记录。"
                        "对方需要先主动给 Bot 发过消息，系统才能记录其 chat_id。"
                    )
            if not chat_id:
                return "错误：read_group_messages 需要提供 chat_id（群）或 p2p_open_id（P2P 私信）参数。"
            import time as _time
            # 优先使用明确的时间段；否则按 hours 回退
            if start_time and end_time:
                import pytz as _pytz
                from datetime import datetime as _datetime
                _TZ = _pytz.timezone("Asia/Shanghai")
                def _parse_dt(s: str) -> int:
                    for fmt in (
                        "%Y-%m-%dT%H:%M:%S%z", "%Y-%m-%dT%H:%M%z",
                        "%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M", "%Y-%m-%d",
                    ):
                        try:
                            dt = _datetime.strptime(s.strip(), fmt)
                            if dt.tzinfo is None:
                                dt = _TZ.localize(dt)
                            return int(dt.timestamp())
                        except ValueError:
                            continue
                    raise ValueError(f"无法解析时间格式: {s!r}")
                try:
                    start_ts = _parse_dt(start_time)
                    end_ts   = _parse_dt(end_time)
                except ValueError as ve:
                    return f"时间格式错误：{ve}。请使用 'YYYY-MM-DD HH:MM' 格式。"
            else:
                end_ts   = int(_time.time())
                start_ts = end_ts - hours * 3600
            msgs = []
            source = "API"
            api_error_note = ""
            try:
                msgs = self.feishu.get_group_messages(chat_id, start_ts, end_ts, max_msgs=300)
                if msgs:
                    source = "API"
            except Exception as e:
                logger.warning("get_group_messages API 失败: %s", e)
            # Bot token 无结果时，显式尝试用户 token（Bot 不在群/曾被移出时）
            if not msgs:
                try:
                    user_msgs = self.feishu.get_group_messages_as_user(
                        chat_id, start_ts, end_ts, max_msgs=300
                    )
                    if user_msgs:
                        msgs = user_msgs
                        source = "用户身份API"
                except Exception as ue:
                    logger.warning("get_group_messages_as_user 失败: %s", ue)
                    ue_str = str(ue)
                    if "231204" in ue_str or "99991679" in ue_str:
                        api_error_note = "（提示：用户授权 token 缺少群消息权限，请重新运行 scripts/authorize_user_im.py 完成授权）"
                    elif "230002" in ue_str:
                        api_error_note = "（提示：Bot 不在该群内，且用户授权 token 未就绪，请运行 scripts/authorize_user_im.py 完成授权）"
                    else:
                        api_error_note = f"（API 错误：{ue}）"
            # 最终降级到本地消息缓存
            if not msgs and self._message_cache:
                msgs = self._message_cache.get_messages(chat_id, start_ts, end_ts, max_msgs=300)
                source = "本地缓存"
            # 用于展示的时间描述
            if start_time and end_time:
                time_desc = f"{start_time} ~ {end_time}"
            else:
                time_desc = f"最近 {hours} 小时"
            if not msgs:
                return f"{time_desc} 内没有消息{api_error_note}。"
            lines = []
            for m in msgs:
                if m.get("is_thread_reply"):
                    lines.append(f"  └ [{m.get('sender_name', '?')}](话题回复) {m.get('text', '')}")
                else:
                    mid = m.get("message_id", "")
                    mid_str = f" [id:{mid}]" if mid else ""
                    lines.append(f"[{m.get('sender_name', '?')}]{mid_str} {m.get('text', '')}")
            return f"{time_desc} 共 {len(msgs)} 条消息（来源：{source}）：\n\n" + "\n".join(lines)
        if action == "apply_mentions":
            if not url:
                return "错误：apply_mentions 需要提供 url 参数（wiki 页面 URL）。"
            if not mention_map:
                return "错误：apply_mentions 需要提供 mention_map 参数（{姓名: open_id}）。"
            # 从 URL 中提取 node_token（wiki URL 格式：.../wiki/{token}）
            import re as _re
            m = _re.search(r"/wiki/(\w+)", url)
            if not m:
                return f"无法从 URL 中提取 wiki node_token: {url}"
            node_token = m.group(1)
            ok, detail = self.feishu.apply_mentions_to_wiki_page(node_token, mention_map)
            return f"@mention 写入{'成功' if ok else '失败'}：{detail}"
        if action == "inspect_pptx":
            if not obj_token:
                return "错误：inspect_pptx 需要提供 obj_token（从 read_page 结果的【文件 token：xxx】中获取）。"
            return self._inspect_pptx(obj_token, slide_index)
        if action == "edit_pptx":
            if not obj_token:
                return "错误：edit_pptx 需要提供 obj_token（从 read_page 结果的【文件 token：xxx】中获取）。"
            if not replacements and not shape_updates:
                return "错误：edit_pptx 需要提供 replacements（文本替换）或 shape_updates（形状调整）中至少一个。"
            return self._edit_pptx_in_place(obj_token, replacements or [], shape_updates or [], output_filename)
        if action == "create_task":
            if not task_title:
                return "错误：create_task 操作需要提供 task_title 参数。"
            # 解析截止时间 → Unix 毫秒
            due_ms: int | None = None
            if task_due:
                import pytz as _pytz
                _tz = _pytz.timezone("Asia/Shanghai")
                _parsed = None
                for _fmt in ("%Y-%m-%dT%H:%M:%S", "%Y-%m-%d %H:%M", "%Y-%m-%d"):
                    try:
                        from datetime import datetime as _dt
                        _parsed = _dt.strptime(task_due, _fmt)
                        break
                    except ValueError:
                        continue
                if _parsed is None:
                    return f"错误：无法解析截止时间 '{task_due}'，请使用格式：'2026-03-21 18:00' 或 '2026-03-21'。"
                due_ms = int(_tz.localize(_parsed).timestamp() * 1000)
            try:
                result = self.feishu.create_task(
                    title=task_title,
                    description=task_description or "",
                    due_ms=due_ms,
                    assignee_open_ids=task_assignee_open_ids or [],
                    follower_open_ids=task_follower_open_ids or [],
                )
                assignee_n = len(task_assignee_open_ids or [])
                follower_n = len(task_follower_open_ids or [])
                due_note = f"\n截止时间：{task_due}" if task_due else ""
                return (
                    f"任务已创建！\n"
                    f"标题：{result['summary']}{due_note}\n"
                    f"负责人：{assignee_n} 人，关注者：{follower_n} 人\n"
                    f"任务链接：{result['url']}\n"
                    f"[task_id={result['task_id']}]"
                )
            except Exception as e:
                return f"创建飞书任务失败：{e}\n提示：请确认应用已开启 task:task:write 权限，并在飞书开放平台重新发布。"
        if action == "create_doc":
            if not title or not content:
                return "错误：create_doc 需要提供 title 和 content（Markdown）参数。"
            script = os.path.join(_FEISHU_DOC_SCRIPTS_DIR, "md2feishu.py")
            if not os.path.exists(script):
                return f"错误：feishu-doc skill 未安装，找不到 {script}"
            try:
                with tempfile.NamedTemporaryFile(
                    mode="w", suffix=".md", delete=False, encoding="utf-8"
                ) as f:
                    f.write(content)
                    tmp_path = f.name
                cmd = [sys.executable, script, "create", tmp_path, "--title", title]
                if ref_url:
                    cmd += ["--wiki", ref_url]
                env = {**os.environ, "FEISHU_APP_ID": self.feishu.app_id, "FEISHU_APP_SECRET": self.feishu.app_secret}
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60, env=env)
                os.unlink(tmp_path)
                output = (result.stdout or result.stderr or "").strip()
                if result.returncode != 0:
                    return f"创建飞书文档失败：\n{output[:2000]}"
                return f"飞书文档已创建！\n{output[:1000]}"
            except Exception as e:
                return f"创建飞书文档失败：{e}"
        if action == "drawio_to_board":
            if not title or not drawio_content:
                return "错误：drawio_to_board 需要提供 title 和 drawio_content（drawio XML）参数。"
            script = os.path.join(_FEISHU_DOC_SCRIPTS_DIR, "drawio2board.py")
            if not os.path.exists(script):
                return f"错误：feishu-doc skill 未安装，找不到 {script}"
            try:
                with tempfile.NamedTemporaryFile(
                    mode="w", suffix=".drawio", delete=False, encoding="utf-8"
                ) as f:
                    f.write(drawio_content)
                    tmp_path = f.name
                cmd = [sys.executable, script, "create", tmp_path, "--title", title]
                if ref_url:
                    cmd += ["--wiki", ref_url]
                env = {**os.environ, "FEISHU_APP_ID": self.feishu.app_id, "FEISHU_APP_SECRET": self.feishu.app_secret}
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60, env=env)
                os.unlink(tmp_path)
                output = (result.stdout or result.stderr or "").strip()
                if result.returncode != 0:
                    return f"drawio 转飞书画板失败：\n{output[:2000]}"
                return f"飞书画板已创建！\n{output[:1000]}"
            except Exception as e:
                return f"drawio 转飞书画板失败：{e}"
        if action == "group_summary":
            if not self._daily_summary:
                return "group_summary 不可用：DailySummaryJob 未注入（请联系管理员）"
            allowed = ("biweekly", "monthly", "quarterly")
            if report_type not in allowed:
                return f"group_summary 需要 report_type，可选值：{allowed}。daily/weekly 由系统自动触发。"
            try:
                results = self._daily_summary.run_for_report_type(report_type, days_back)
                if results:
                    return f"{report_type} 摘要已完成，共处理 {len(results)} 个群，摘要已发布到飞书 Wiki 并发送 DM 通知。"
                return f"{report_type} 摘要完成：在指定时间段内未找到符合条件的群消息（tier 范围内无活跃群）。"
            except Exception as e:
                return f"group_summary 执行失败：{e}"

        if action == "calendar_events":
            try:
                raw = self._run_feishu_cli("calendar_events", timeout=20)
                return raw if raw else "未获取到日历事件"
            except Exception as e:
                return f"查询日历失败：{e}"

        return f"未知 feishu_action: {action}。支持: search, read_page, list_pages, create_page, edit_page, move_page, send_message, recall_message, list_groups, create_group, get_group_members, read_group_messages, apply_mentions, inspect_pptx, edit_pptx, create_task, create_doc, drawio_to_board, group_summary, calendar_events"

    def _run_feishu_cli(self, *args: str, timeout: int = 30) -> str:
        """调用 feishu-sync-cli，返回 stdout 字符串；失败时自动重试最多 2 次"""
        import time as _time
        cmd = [*_FEISHU_SYNC_CMD, *args]
        logger.info("执行 feishu-sync-cli: %s", " ".join(cmd))
        last_err = ""
        for attempt in range(3):  # 最多尝试 3 次（0, 1, 2）
            try:
                # 用全局锁串行化，防止多线程并发读写同一份 token 文件产生竞态
                with _FEISHU_CLI_LOCK:
                    env = os.environ.copy()
                    env["PYTHONIOENCODING"] = "utf-8"
                    result = subprocess.run(
                        cmd,
                        capture_output=True,
                        text=True,
                        timeout=timeout,
                        encoding="utf-8",
                        errors="replace",
                        env=env,
                    )
                if result.returncode == 0:
                    return result.stdout.strip()

                # stdout 含实际错误，stderr 含 INFO 日志；优先 stdout
                err = (result.stdout or result.stderr or "").strip()
                logger.warning("feishu-sync-cli 返回非零（第%d次）: %s", attempt + 1, err[:1000])

                # 识别已知静态错误（无需重试）
                if "Document type 'file' is not supported" in err or \
                        ("NotImplementedError" in err and "file" in err):
                    return "[飞书工具错误] Document type 'file' is not supported by Feishu API."

                last_err = err
                # token 刷新崩溃或未知错误 → 等待后重试
                if attempt < 2:
                    logger.info("feishu-sync-cli 失败，3 秒后重试（第%d次）...", attempt + 2)
                    _time.sleep(3)

            except FileNotFoundError:
                return (
                    "[飞书工具未安装] feishu-sync-cli 未找到。"
                    "请运行：python3 -m pip install --extra-index-url "
                    "https://artifactory.momenta.works/artifactory/api/pypi/pypi-momenta/simple "
                    "feishu-sync -U"
                )
            except subprocess.TimeoutExpired:
                return f"[超时] feishu-sync-cli 执行超时（>{timeout}s）"
            except Exception as e:
                return f"[错误] feishu-sync-cli 调用失败: {e}"

        return f"[飞书工具错误] {last_err[:500]}"

    def _search_feishu(self, query: str, limit: int = 5) -> str:
        out = self._run_feishu_cli("search_wiki", query, timeout=30)
        if out.startswith("[飞书"):
            return out  # error message from _run_feishu_cli
        if not out:
            return f"未找到与 '{query}' 相关的飞书页面。"
        # 尝试解析 JSON 数组
        try:
            items = json.loads(out)
            if not items:
                return f"未找到与 '{query}' 相关的飞书页面。"
            items = items[:limit]
            lines = [f"找到 {len(items)} 条相关页面："]
            for it in items:
                lines.append(f"- {it.get('title', '(无标题)')}  {it.get('url', '')}")
            return "\n".join(lines)
        except (json.JSONDecodeError, TypeError):
            # 返回原始文本（可能是旧版格式）
            return out[:3000]

    def _read_feishu_page(self, url: str) -> str:
        out = self._run_feishu_cli("read_page_as_markdown", url, timeout=60)
        if out.startswith("["):
            # feishu-sync-cli 失败 → 无论何种原因，先尝试直接下载解析（兼容 PPTX/PDF 及其他 file 类型）
            logger.warning("read_page_as_markdown 失败（%s），尝试 fallback 下载解析", out[:200])
            fallback = self._download_and_parse_wiki_file(url)
            if not fallback.startswith("[") and not fallback.startswith("错误"):
                return fallback
            # fallback 也失败，返回原始错误
            return out
        # 截断过长内容（防止 token 超限）
        if len(out) > 8000:
            out = out[:8000] + "\n\n[...内容过长，已截断，如需完整内容请缩小范围]"
        return out or "页面内容为空。"

    def _download_and_parse_wiki_file(self, url: str) -> str:
        """处理 feishu wiki 中 obj_type=file 的节点（如 PPTX/PDF）：下载并提取文字。"""
        import re as _re
        import glob as _glob
        import json as _json

        # 从 URL 中提取 node_token
        m = _re.search(r"/wiki/(\w+)", url)
        if not m:
            return "[文件类型不支持] 该页面是二进制文件，无法通过文档 API 读取。"
        node_token = m.group(1)

        # 从 feishu-sync 本地缓存读取 obj_token（feishu-sync-cli 首次访问时会写入 .info.json）
        pattern = os.path.expanduser(
            f"~/.cache/feishu-sync/export/wiki/*/*/{node_token}.info.json"
        )
        matches = _glob.glob(pattern)
        if not matches:
            return (
                "[文件类型不支持] 该页面是二进制文件（如 PPTX/PDF），无法通过文档 API 读取。"
                "请将文件内容粘贴为文字，或以飞书文档形式提供。"
            )

        info = _json.load(open(matches[0]))
        obj_token = info.get("obj_token", "")
        title = info.get("title", "未知文件")

        if not obj_token:
            return f"[文件类型不支持] 未找到 {title} 的下载 token。"

        # 下载文件
        logger.info("下载 wiki 文件: %s (obj_token=%s)", title, obj_token)
        data = self.feishu.download_drive_file(obj_token)
        if not data:
            return f"[文件下载失败] 无法下载 {title}，请检查 Bot 是否有该文件的读取权限。"

        # 按文件类型解析
        lower_title = title.lower()
        if lower_title.endswith(".pptx"):
            try:
                text = self.feishu.parse_pptx_text(data)
                slide_count = text.count("### 第")
                out = f"## {title}\n\n（PPTX，共 {slide_count} 页）\n\n{text}"
                if len(out) > 12000:
                    out = out[:12000] + "\n\n[...内容过长，已截断]"
                out += f"\n\n【文件 token：{obj_token}（如需在原文件上修改，请使用 feishu_action edit_pptx + 此 token）】"
                logger.info("PPTX 解析成功: %s，%d 页", title, slide_count)
                return out
            except Exception as e:
                return f"[PPTX 解析失败] {title}: {e}"
        else:
            ext = title.rsplit(".", 1)[-1].upper() if "." in title else "未知"
            return (
                f"[不支持的文件格式] {title} 是 {ext} 格式。"
                "目前支持解析 PPTX 文件；PDF、DOCX 等格式如需支持请告知。"
            )

    def _load_pptx_from_feishu(self, obj_token: str):
        """下载飞书云盘 PPTX 并返回 (Presentation, bytes)，失败返回 None。"""
        import io as _io
        try:
            from pptx import Presentation as _Prs
        except ImportError:
            import subprocess as _sp
            _sp.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
            from pptx import Presentation as _Prs
        data = self.feishu.download_drive_file(obj_token)
        if not data:
            return None, None
        return _Prs(_io.BytesIO(data)), data

    def _inspect_pptx(self, obj_token: str, slide_index: int | None = None) -> str:
        """列出 PPTX 每张幻灯片中所有形状的坐标、尺寸和文字预览，便于后续精确修改形状位置。"""
        from pptx.util import Inches
        prs, _ = self._load_pptx_from_feishu(obj_token)
        if prs is None:
            return "PPTX 下载失败，请确认 obj_token 是否正确。"

        lines = [f"PPTX 共 {len(prs.slides)} 页\n"]
        slide_iter = (
            [(slide_index, prs.slides[slide_index])]
            if slide_index is not None
            else list(enumerate(prs.slides))
        )
        for s_idx, slide in slide_iter:
            lines.append(f"### 第 {s_idx + 1} 页（slide_index={s_idx}）")
            for sh_idx, shape in enumerate(slide.shapes):
                left   = round(shape.left   / 914400, 3)  # EMU → 英寸
                top    = round(shape.top    / 914400, 3)
                width  = round(shape.width  / 914400, 3)
                height = round(shape.height / 914400, 3)
                text_preview = ""
                if hasattr(shape, "text_frame"):
                    text_preview = shape.text_frame.text[:40].replace("\n", "↵")
                elif hasattr(shape, "text"):
                    text_preview = str(shape.text)[:40]
                shape_type = shape.shape_type  # 1=rect,2=line,6=picture,13=textbox...
                lines.append(
                    f"  shape_index={sh_idx}  name={shape.name!r}  type={shape_type}"
                    f"  left={left}\" top={top}\" w={width}\" h={height}\""
                    + (f"  text={text_preview!r}" if text_preview else "")
                )
        out = "\n".join(lines)
        if len(out) > 5000:
            out = out[:5000] + "\n\n[...已截断，如需查看更多形状请指定 slide_index]"
        out += "\n\n⚠️ 以上是形状坐标信息，仅供参考。请立即调用 edit_pptx 应用实际修改（replacements + shape_updates），否则不会有任何变化。"
        return out

    def _edit_pptx_in_place(
        self,
        obj_token: str,
        replacements: list[dict],
        shape_updates: list[dict] | None = None,
        output_filename: str = "修改后.pptx",
    ) -> str:
        """下载飞书云盘 PPTX，做文本替换 + 形状位置/尺寸调整（保留 run 格式），保存到 PPT_OUTPUT_DIR。"""
        prs, _ = self._load_pptx_from_feishu(obj_token)
        if prs is None:
            return "PPTX 下载失败，请确认 obj_token 是否正确，以及 Bot 是否有该文件的读取权限。"

        # ── 文本替换 ────────────────────────────────────────────────
        text_changed = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for rep in (replacements or []):
                            old = rep.get("old", "")
                            new = rep.get("new", "")
                            if old and old in run.text:
                                run.text = run.text.replace(old, new)
                                text_changed += 1

        # ── 形状位置/尺寸调整 ────────────────────────────────────────
        shape_changed = 0
        EMU = 914400  # 1 inch in EMU
        for upd in (shape_updates or []):
            s_idx = upd.get("slide")
            sh_idx = upd.get("shape")
            if s_idx is None or sh_idx is None:
                continue
            try:
                shape = prs.slides[s_idx].shapes[sh_idx]
            except IndexError:
                continue
            if "left"   in upd: shape.left   = int(upd["left"]   * EMU)
            if "top"    in upd: shape.top    = int(upd["top"]    * EMU)
            if "width"  in upd: shape.width  = int(upd["width"]  * EMU)
            if "height" in upd: shape.height = int(upd["height"] * EMU)
            shape_changed += 1

        # 重新读取（运行时 env var 可能比模块加载时晚，所以每次从 env 读）
        out_dir = os.environ.get("MOMENTA_PPT_DIR") or os.path.join(
            os.path.expanduser("~"), "Desktop", "Claude"
        )
        os.makedirs(out_dir, exist_ok=True)
        if not output_filename.endswith(".pptx"):
            output_filename += ".pptx"
        output_path = os.path.join(out_dir, output_filename)
        prs.save(output_path)
        logger.info("PPTX 已修改保存: %s，文本替换 %d 处，形状调整 %d 个", output_path, text_changed, shape_changed)

        parts = []
        if text_changed:
            parts.append(f"替换了 {text_changed} 处文本")
        if shape_changed:
            parts.append(f"调整了 {shape_changed} 个形状位置/尺寸")
        summary = "、".join(parts) or "无变更"

        # 上传到飞书 Drive（与原文件同目录）
        drive_url = ""
        try:
            parent_token, parent_type = self.feishu.get_drive_file_parent(obj_token)
            if parent_token:
                with open(output_path, "rb") as f:
                    file_bytes = f.read()
                _, drive_url = self.feishu.upload_drive_file(
                    output_filename, file_bytes, parent_token, parent_type
                )
        except Exception as e:
            logger.warning("Drive 上传失败（不影响本地保存）: %s", e)

        result = (
            f"PPTX 已修改！\n"
            f"修改内容：{summary}（保留原始格式）\n"
            f"本地路径：{output_path}"
        )
        if drive_url:
            result += f"\n飞书链接：{drive_url}"
        else:
            result += "\n（飞书 Drive 上传失败，请手动上传本地文件）"
        return result

    def _list_feishu_pages(self, url: str) -> str:
        out = self._run_feishu_cli("list_pages", url, timeout=30)
        if out.startswith("["):
            return out
        return out or "该空间/页面下没有子页面。"

    def _create_feishu_page(
        self, ref_url: str, title: str, content: str = "", position: str = "child"
    ) -> str:
        args = ["create_page", ref_url, title]
        if content:
            args.append(content)
        if position == "sibling":
            args.append("--position=sibling")
        out = self._run_feishu_cli(*args, timeout=30)
        if out.startswith("["):
            return out
        return out or f"页面 '{title}' 已创建。"

    def _edit_feishu_page(self, url: str, old_string: str, new_string: str) -> str:
        out = self._run_feishu_cli(
            "edit_page", url,
            f"--old_string={old_string}",
            f"--new_string={new_string}",
            timeout=30,
        )
        if out.startswith("["):
            return out
        return out or "页面已更新。"

    def _move_feishu_page(self, url: str, target_url: str) -> str:
        """将 Wiki 页面移动到目标父目录下"""
        # 从 URL 末尾提取 node_token（wiki URL 格式：.../wiki/<token>）
        node_token = url.rstrip("/").split("/")[-1]
        target_parent_token = target_url.rstrip("/").split("/")[-1]

        # 直接通过飞书 API 获取 space_id（避免依赖可能过期的 CLI 缓存）
        space_id = self.feishu.get_wiki_node_space_id(node_token)
        if not space_id:
            return "[错误] 无法获取页面所在空间 ID，请确认该页面是 Wiki 页面，且应用有读取权限。"

        result = self.feishu.move_wiki_page(space_id, node_token, target_parent_token)
        if not result:
            return "移动失败，请检查目标目录 URL 是否正确，以及应用是否有编辑权限。"

        new_url = f"https://momenta.feishu.cn/wiki/{result.get('node_token', node_token)}"
        return f"页面已移动成功！\n新位置：{new_url}"

    # ── Bag / Mviz 分析 ────────────────────────────────────────────

    def _bag_analysis(
        self,
        action: str,
        url: str = "",
        topic: str = "",
        stream_name: str = "",
        query: str = "",
        bag_md5: str = "",
        max_frames: int = 15,
    ) -> str:
        """统一分发 Bag/Mviz 分析操作"""
        if action == "search_streams":
            if not query:
                return "错误：search_streams 需要提供 query 参数。"
            return self._search_mviz_streams(query)
        if action == "lookup_stream":
            if not stream_name:
                return "错误：lookup_stream 需要提供 stream_name 参数。"
            return self._lookup_mviz_stream(stream_name)
        if action == "resolve_url":
            if not url:
                return "错误：resolve_url 需要提供 url 参数。"
            result = self._do_resolve_mviz_url(url)
            if isinstance(result, str):
                return result
            lines = [
                "✅ URL 解析成功",
                f"环境: {result['env']}",
                f"bag_md5: {result['md5'] or '（无）'}",
                f"storage_name: {result['storage_name'] or '（无）'}",
                f"storage_key: {(result['storage_key'] or '')[:60] or '（无）'}",
                f"API 端点: {result['api_endpoint']}",
            ]
            return "\n".join(lines)
        if action == "download_topic":
            if not topic:
                return "错误：download_topic 需要提供 topic 参数。"
            if not url and not bag_md5:
                return "错误：download_topic 需要提供 url 或 bag_md5 参数。"
            return self._download_bag_topic(url, topic, max_frames, bag_md5=bag_md5)
        return f"未知 bag_analysis action: {action}，支持: search_streams, lookup_stream, resolve_url, download_topic"

    def _search_mviz_streams(self, query: str) -> str:
        """在本地 stream-catalog.json 中按关键词搜索 stream"""
        catalog = _load_stream_catalog()
        if isinstance(catalog, str):
            return catalog  # 错误信息

        streams = catalog.get("streams", [])
        q = query.lower()
        matches = []
        for s in streams:
            name = s.get("streamName", "")
            desc = s.get("description", "")
            topics = " ".join(s.get("topics", []))
            formatter = s.get("formatter", "")
            if q in name.lower() or q in desc.lower() or q in topics.lower() or q in formatter.lower():
                matches.append(s)
        if not matches:
            return f"未找到与 '{query}' 相关的 stream。尝试使用英文关键词或 topic 名称。"
        lines = [f"找到 {len(matches)} 个相关 stream（显示前 20）："]
        for s in matches[:20]:
            topics = ", ".join(s.get("topics", []))
            lines.append(f"  • {s.get('streamName', '')}  ← {topics or '(无 topic)'}")
        return "\n".join(lines)

    def _lookup_mviz_stream(self, stream_name: str) -> str:
        """精确查询某 stream 的详细信息"""
        catalog = _load_stream_catalog()
        if isinstance(catalog, str):
            return catalog  # 错误信息

        streams = catalog.get("streams", [])
        # 精确匹配
        match = next((s for s in streams if s.get("streamName") == stream_name), None)
        # 如果精确匹配失败，做模糊搜索并给出候选
        if not match:
            q = stream_name.lower()
            candidates = [s for s in streams if q in s.get("streamName", "").lower()][:5]
            if candidates:
                cands = "\n".join(f"  • {c['streamName']}" for c in candidates)
                return f"未找到精确 stream '{stream_name}'。相似候选：\n{cands}"
            return f"未找到 stream '{stream_name}'，请使用 search_streams 查找。"

        lines = [
            f"Stream: {match.get('streamName', '')}",
            f"描述: {match.get('description', '（无）')}",
            f"",
            f"订阅 Topics:",
        ]
        for t in match.get("topics", []):
            lines.append(f"  • {t}")
        if match.get("formatter"):
            lines.append(f"Formatter: {match['formatter']}")
        if match.get("parsers"):
            lines.append(f"Parsers:")
            for p in match["parsers"][:5]:
                lines.append(f"  • {p}")
        if match.get("coordinateSystem"):
            lines.append(f"坐标系: {match['coordinateSystem']}")
        if match.get("renderType"):
            lines.append(f"渲染类型: {match['renderType']}")
        return "\n".join(lines)

    def _do_resolve_mviz_url(self, url: str) -> dict | str:
        """解析 Mviz URL，返回 {md5, storage_name, storage_key, api_endpoint, env} 或错误字符串"""
        # 命中缓存时直接返回，避免重复 HTTP 请求
        if url in self._url_resolve_cache:
            logger.debug("URL 解析缓存命中: %s", url[:80])
            return self._url_resolve_cache[url]

        # 检测环境
        if "cla.eu" in url:
            env = "EU"
        elif ".jp." in url or "jp.momenta" in url:
            env = "JP"
        elif ".me." in url or "me.momenta" in url:
            env = "ME"
        else:
            env = "CN"

        api_endpoint = _TOPIC_ECHO_ENDPOINTS[env]
        meta_api_domain = _META_API_DOMAINS[env]

        md5 = ""
        storage_name = ""
        storage_key = ""

        # 尝试提取 bag_md5
        m = re.search(r"bag_md5=([a-f0-9]{32})", url)
        if not m:
            m = re.search(r"bag_md5=([^&\s]+)", url)
        if m:
            md5 = m.group(1)

        # 如果没有 md5，尝试 meta 参数
        if not md5:
            m = re.search(r"meta=([^&\s]+)", url)
            if m:
                meta_param = m.group(1)
                # Step 1: Meta API → bucket + bag_key
                meta_api_url = f"https://{meta_api_domain}/api/v1/web/proxy/get?meta={meta_param}"
                try:
                    resp = _requests.get(meta_api_url, headers={"accept": "*/*"}, timeout=15)
                    resp.raise_for_status()
                    meta_data = resp.json()
                except Exception as e:
                    return f"[错误] Meta API 调用失败: {e}"
                bucket = meta_data.get("bucket", "")
                bag_key = meta_data.get("bag_key", "")
                if not bucket or not bag_key:
                    return f"[错误] Meta API 返回数据缺少字段：{meta_data}"

                # Step 2: Bag Result API → storage_name + md5 + storage_key
                encoded_key = urllib.parse.quote(bag_key)
                bag_result_url = (
                    f"https://{meta_api_domain}/api/bag/result"
                    f"?bucket={bucket}&bag_key={encoded_key}&scene=label&from_prod=true"
                )
                try:
                    resp = _requests.get(bag_result_url, headers={"accept": "*/*"}, timeout=15)
                    resp.raise_for_status()
                    result_data = resp.json()
                except Exception as e:
                    return f"[错误] Bag Result API 调用失败: {e}"
                d = result_data.get("data", {})
                storage_name = d.get("storage_name", "")
                md5 = d.get("bag_md5", "")
                storage_key = d.get("bag_key", "")
                if not md5:
                    return f"[错误] 无法获取 bag_md5：{result_data}"

        if not md5 and not storage_name:
            return "[错误] 无法从 URL 提取 bag 参数，需要包含 bag_md5= 或 meta= 参数"

        result = {
            "md5": md5,
            "storage_name": storage_name,
            "storage_key": storage_key,
            "api_endpoint": api_endpoint,
            "env": env,
        }
        self._url_resolve_cache[url] = result
        return result

    def _download_bag_topic(
        self, url: str, topic: str, max_frames: int = 15, *, bag_md5: str = ""
    ) -> str:
        """从 Mviz URL/MD5 下载指定 topic 数据并返回分析摘要"""
        # 1. 解析 URL（若传入 bag_md5 且 URL 已缓存，可快速获取 storage 信息）
        if bag_md5 and not url:
            # 纯 md5 模式，默认 CN 环境，无 storage 信息
            params: dict | str = {
                "md5": bag_md5,
                "storage_name": "",
                "storage_key": "",
                "api_endpoint": _TOPIC_ECHO_ENDPOINTS["CN"],
                "env": "CN",
            }
        else:
            params = self._do_resolve_mviz_url(url)
        if isinstance(params, str):
            return params  # 错误信息

        md5 = bag_md5 or params["md5"]
        storage_name = params["storage_name"]
        storage_key = params["storage_key"]
        api_endpoint = params["api_endpoint"]

        logger.info("下载 topic %s from bag %s (env=%s)", topic, md5, params["env"])

        # 2. 构建下载 URL
        encoded_topic = urllib.parse.quote(topic)
        only_md5 = not storage_name or storage_name in _ONLY_MD5_STORAGES or storage_name.endswith("-dxp")

        if only_md5 and md5:
            download_url = f"{api_endpoint}?md5={md5}&topic={encoded_topic}&convert_proto=1"
        elif storage_name and storage_key:
            encoded_sk = urllib.parse.quote(storage_key)
            download_url = (
                f"{api_endpoint}?md5={md5}&storage_name={storage_name}"
                f"&storage_key={encoded_sk}&topic={encoded_topic}&convert_proto=1"
            )
        else:
            download_url = f"{api_endpoint}?md5={md5}&topic={encoded_topic}&convert_proto=1"

        # 3. 下载数据
        try:
            resp = _requests.get(download_url, headers={"accept": "*/*"}, timeout=60)
            resp.raise_for_status()
        except _requests.exceptions.HTTPError as e:
            return f"[HTTP 错误] {e}\n下载 URL：{download_url}"
        except Exception as e:
            return f"[下载错误] {e}\n下载 URL：{download_url}"

        # 4. 解析 JSON
        try:
            data = resp.json()
        except Exception:
            snippet = resp.text[:1000]
            return f"[警告] 响应非 JSON（可能是 protobuf 或空数据）:\n{snippet}"

        # 5. 格式化输出
        return self._format_topic_data(topic, data, max_frames, md5)

    def _format_topic_data(self, topic: str, data, max_frames: int, bag_md5: str) -> str:
        """将 topic 数据格式化为可读摘要"""
        # 兼容多种响应结构
        if isinstance(data, list):
            frames = data
        elif isinstance(data, dict):
            # {"data": [...]} 或 {"messages": [...]}
            frames = data.get("data") or data.get("messages") or [data]
        else:
            return f"[未知数据格式] type={type(data)}"

        if not frames:
            return (
                f"Topic: {topic}\nbag: {bag_md5}\n"
                f"⚠️ 没有数据（该 bag 中可能不存在此 topic，或 topic 名称有误）"
            )

        total = len(frames)
        sample = frames[:max_frames]

        lines = [
            f"Topic: {topic}",
            f"Bag: {bag_md5}",
            f"总帧数: {total}，显示前 {len(sample)} 帧",
            f"",
        ]

        for i, frame in enumerate(sample):
            # 递归转换 {sec, nsec} 时间戳为浮点数
            frame = _convert_timestamps(frame)
            frame_str = json.dumps(frame, ensure_ascii=False, indent=None)
            # 截断单帧
            if len(frame_str) > 600:
                frame_str = frame_str[:600] + " ..."
            lines.append(f"[{i+1}] {frame_str}")

        if total > max_frames:
            lines.append(f"\n... 还有 {total - max_frames} 帧未显示 ...")

        return "\n".join(lines)

    # ── PPT 生成 ──────────────────────────────────────────────────

    def _generate_ppt(
        self,
        title: str,
        slides: list[dict],
        subtitle: str = "",
        output_filename: str = "客户汇报.pptx",
    ) -> str:
        # 每次调用时从 env 读，确保 .env 晚加载也能生效
        out_dir = os.environ.get("MOMENTA_PPT_DIR") or os.path.join(
            os.path.expanduser("~"), "Desktop", "Claude"
        )
        os.makedirs(out_dir, exist_ok=True)

        # 确保文件名以 .pptx 结尾
        if not output_filename.endswith(".pptx"):
            output_filename += ".pptx"
        output_path = os.path.join(out_dir, output_filename)

        # 将数据序列化后传给生成脚本
        ppt_data = {
            "title": title,
            "subtitle": subtitle,
            "slides": slides,
            "output_path": output_path,
            "template_path": PPT_TEMPLATE_PATH if os.path.exists(PPT_TEMPLATE_PATH) else "",
        }

        script = _build_ppt_script(ppt_data)

        # 写入临时文件并执行
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".py", delete=False, encoding="utf-8"
        ) as f:
            f.write(script)
            tmp_path = f.name

        try:
            result = subprocess.run(
                [sys.executable, tmp_path],
                capture_output=True,
                text=True,
                timeout=60,
                encoding="utf-8",
                errors="replace",
            )
            if result.returncode != 0:
                err = (result.stderr or result.stdout or "").strip()
                logger.error("PPT 生成失败: %s", err[:500])
                return f"PPT 生成失败：{err[:400]}"
            logger.info("PPT 已生成: %s", output_path)
            tpl_info = "（使用 Momenta 模版）" if os.path.exists(PPT_TEMPLATE_PATH) else ""

            # 上传到飞书 Drive「我的空间」根目录
            drive_url = ""
            try:
                parent_token, parent_type = self.feishu.get_drive_file_parent("")  # 空 token → fallback 根目录
                if parent_token:
                    with open(output_path, "rb") as f:
                        file_bytes = f.read()
                    _, drive_url = self.feishu.upload_drive_file(
                        output_filename, file_bytes, parent_token, parent_type
                    )
            except Exception as e:
                logger.warning("生成 PPT 后 Drive 上传失败（不影响本地保存）: %s", e)

            msg = (
                f"PPT 已生成{tpl_info}！\n"
                f"共 {1 + len(slides)} 页（封面 + {len(slides)} 内容页）\n"
                f"本地路径：{output_path}"
            )
            if drive_url:
                msg += f"\n飞书链接：{drive_url}"
            return msg
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    # ── PM 项目计划生成 ─────────────────────────────────────────────

    def _generate_pm_plan(
        self,
        project_name: str,
        t0_date: str,
        total_weeks: int = 20,
        drive_folder_token: str = "",
        wiki_parent_url: str = "",
        milestones: list = None,
        modules_override: list = None,
    ) -> str:
        """生成 L2++ 适配项目二级开发计划甘特图（HTML + PNG）"""
        from pathlib import Path

        # ── 生成唯一 slug ──────────────────────────────────────────
        slug_base = re.sub(r"[^a-zA-Z0-9]", "", project_name)[:6] or "pm"
        slug = f"{slug_base}{int(datetime.now().timestamp()) % 10000}"
        workdir = Path(f"/tmp/pm_{slug}")
        workdir.mkdir(parents=True, exist_ok=True)

        # ── 默认里程碑 ─────────────────────────────────────────────
        if not milestones:
            milestones = [
                {"week": 0,            "name": "EP车辆Ready", "version": "首车调试版本"},
                {"week": 4,            "name": "首车",         "version": "CP/APA闭环版本"},
                {"week": 10,           "name": "CMA-base",     "version": "公告送测"},
                {"week": 18,           "name": "终版送测",     "version": "终版送测（性能保证）"},
                {"week": total_weeks,  "name": "SOP版本",      "version": "封版"},
            ]

        # ── 默认模块（标准 L2++ 结构）──────────────────────────────
        modules = modules_override or [
            {"name": "底盘联调", "tasks": [
                {"label": "执行器自标定",        "start": -2, "end": 0,  "color": "#5a6e1f"},
                {"label": "执行器开环验收",       "start": 0,  "end": 4,  "color": "#5a6e1f"},
                {"label": "闭环调整（控制精调）", "start": 4,  "end": 10, "color": "#5a6e1f"},
            ]},
            {"name": "控制匹配", "tasks": [
                {"label": "临时控制参数",   "start": 0,  "end": 4,  "color": "#5a6e1f"},
                {"label": "第一版控制参数", "start": 4,  "end": 6,  "color": "#5a6e1f"},
                {"label": "控制锁参",       "start": 6,  "end": 10, "color": "#5a6e1f"},
            ]},
            {"name": "相机", "tasks": [
                {"label": "V1-基础ISP",   "start": 0, "end": 4,  "color": "#5a6e1f"},
                {"label": "V2-性能版本",  "start": 4, "end": 10, "color": "#5a6e1f"},
            ]},
            {"name": "Radar", "tasks": [
                {"label": "V1-替代参数+通讯联调", "start": 0, "end": 4,  "color": "#5a6e1f"},
                {"label": "V2-首轮验收",           "start": 4, "end": 10, "color": "#5a6e1f"},
            ]},
            {"name": "USS", "tasks": [
                {"label": "V1-功能版本",  "start": 0, "end": 4,  "color": "#5a6e1f"},
                {"label": "V2-USSOD验收", "start": 4, "end": 10, "color": "#5a6e1f"},
            ]},
            {"name": "Lidar", "tasks": [
                {"label": "V1-基础功能+通讯联调", "start": 0, "end": 4,  "color": "#5a6e1f"},
                {"label": "V2-首轮验收",           "start": 4, "end": 10, "color": "#5a6e1f"},
            ]},
            {"name": "域控", "tasks": [
                {"label": "V1-首版底软",   "start": 0, "end": 4,  "color": "#5a6e1f"},
                {"label": "V2-底软稳定版", "start": 4, "end": 10, "color": "#5a6e1f"},
            ]},
            {"name": "CP主线对齐", "tasks": [
                {"label": "对齐主线-首轮（CP）", "start": 10, "end": 14, "color": "#2f855a"},
                {"label": "主线对齐-二轮（CP）", "start": 14, "end": 18, "color": "#2f855a"},
            ]},
            {"name": "AEB/FCW", "tasks": [
                {"label": "标定",     "start": 10, "end": 14, "color": "#2f855a"},
                {"label": "调参",     "start": 14, "end": 16, "color": "#2f855a"},
                {"label": "场地测试", "start": 16, "end": 18, "color": "#2f855a"},
            ]},
            {"name": "APA/RPA", "tasks": [
                {"label": "CMA摸底",      "start": 10, "end": 14, "color": "#2f855a"},
                {"label": "场景构建测试", "start": 14, "end": 16, "color": "#2f855a"},
                {"label": "特性评测",     "start": 16, "end": 18, "color": "#2f855a"},
            ]},
            {"name": "NOA", "tasks": [
                {"label": "NOA（CMA）", "start": 10, "end": 14, "color": "#2f855a"},
                {"label": "泛化",       "start": 14, "end": 18, "color": "#2f855a"},
            ]},
            {"name": "压测/性能优化", "tasks": [
                {"label": "压测、性能优化", "start": 18, "end": total_weeks, "color": "#2f855a"},
            ]},
        ]

        # ── 写 config.json ─────────────────────────────────────────
        config = {
            "project": project_name,
            "t0": t0_date,
            "total_weeks": total_weeks,
            "milestones": milestones,
            "modules": modules,
        }
        (workdir / "config.json").write_text(
            json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8"
        )

        # ── 复制 gen_v2_template.py 并替换 SLUG ────────────────────
        skill_dir = Path(__file__).parent.parent / ".agents/skills/pm"
        template_path = skill_dir / "gen_v2_template.py"
        if not template_path.exists():
            return "错误：pm skill 未安装，找不到 gen_v2_template.py"
        gen_script = template_path.read_text(encoding="utf-8").replace(
            'SLUG = "IS4GR"', f'SLUG = "{slug}"'
        )
        gen_path = workdir / "gen_v2.py"
        gen_path.write_text(gen_script, encoding="utf-8")

        # ── 生成 HTML ──────────────────────────────────────────────
        r = subprocess.run(
            [sys.executable, str(gen_path)],
            capture_output=True, text=True, timeout=30, encoding="utf-8",
        )
        if r.returncode != 0:
            err = (r.stderr or r.stdout or "").strip()
            return f"生成 HTML 甘特图失败：{err[:400]}"
        html_path = workdir / "project_plan.html"
        logger.info("PM 甘特图已生成: %s", html_path)

        # ── puppeteer 截图 PNG（可选）──────────────────────────────
        png_path: Path | None = None
        try:
            node_js = (
                f"const p = require('puppeteer');\n"
                f"(async () => {{\n"
                f"  const b = await p.launch({{headless:true, args:['--no-sandbox']}});\n"
                f"  const pg = await b.newPage();\n"
                f"  await pg.setViewport({{width:1400, height:900}});\n"
                f"  await pg.goto('file://{html_path}', {{waitUntil:'networkidle0'}});\n"
                f"  const h = await pg.evaluate(() => document.body.scrollHeight);\n"
                f"  await pg.setViewport({{width:1400, height:h}});\n"
                f"  await pg.goto('file://{html_path}', {{waitUntil:'networkidle0'}});\n"
                f"  await pg.screenshot({{path:'{workdir}/gantt_hd.png', fullPage:true, type:'png'}});\n"
                f"  await b.close();\n"
                f"}})();\n"
            )
            shot_js = workdir / "screenshot.js"
            shot_js.write_text(node_js, encoding="utf-8")
            sr = subprocess.run(
                ["node", str(shot_js)],
                capture_output=True, text=True, timeout=60, cwd="/tmp",
                env={**os.environ, "NODE_PATH": "/tmp/node_modules"},
            )
            candidate = workdir / "gantt_hd.png"
            if sr.returncode == 0 and candidate.exists():
                png_path = candidate
            else:
                logger.warning("puppeteer 截图失败（忽略）: %s", (sr.stderr or "")[:200])
        except Exception as exc:
            logger.warning("puppeteer 截图异常（忽略）: %s", exc)

        # ── 上传到 Drive ───────────────────────────────────────────
        html_url = png_url = ""
        if drive_folder_token:
            try:
                _, html_url = self.feishu.upload_drive_file(
                    f"{project_name}甘特图.html",
                    html_path.read_bytes(), drive_folder_token, "explorer",
                )
                if png_path:
                    _, png_url = self.feishu.upload_drive_file(
                        f"{project_name}甘特图.png",
                        png_path.read_bytes(), drive_folder_token, "explorer",
                    )
            except Exception as exc:
                logger.warning("Drive 上传失败（忽略）: %s", exc)

        # ── 创建 Wiki 页面 ─────────────────────────────────────────
        wiki_url = ""
        if wiki_parent_url:
            try:
                lines = [
                    f"# {project_name} 二级开发计划 V1.0\n",
                    f"> T0（EP车辆Ready）：{t0_date}  |  总周期：{total_weeks} 周\n\n",
                    "## 一、里程碑时间线\n\n",
                    "| 时间点 | 里程碑 | 绝对日期 | 版本计划 |\n",
                    "|--------|--------|----------|----------|\n",
                ]
                for m in milestones:
                    abs_date = (datetime.strptime(t0_date, "%Y-%m-%d") + timedelta(weeks=m["week"])).strftime("%Y-%m-%d")
                    tag = "T0" if m["week"] == 0 else f"T+{m['week']}W"
                    lines.append(f"| {tag} | {m['name']} | {abs_date} | {m.get('version','')} |\n")
                lines.append("\n## 二、各模块开发计划\n\n")
                for mod in modules:
                    lines.append(f"### {mod['name']}\n\n")
                    for t in mod["tasks"]:
                        s = "T前" if t["start"] < 0 else ("T0" if t["start"] == 0 else f"T+{t['start']}W")
                        lines.append(f"- **{t['label']}**（{s} ~ T+{t['end']}W）\n")
                    lines.append("\n")
                lines += [
                    "## 三、关键依赖\n\n",
                    "1. **EP车辆需求**：车辆与锁定车身参数一致；仪表盘无故障\n",
                    "2. **对手件需求**：传感器 +4W 验收软件；+10W 性能锁定\n",
                    "3. **座舱**：公告测试软件提前沟通联调计划及基线\n",
                ]
                md_path = workdir / "feishu_summary.md"
                md_path.write_text("".join(lines), encoding="utf-8")

                feishu_sync = os.path.expanduser("~/Library/Python/3.9/bin/feishu-sync-cli")
                cr = subprocess.run(
                    [feishu_sync, "create_page", wiki_parent_url,
                     f"{project_name} 二级开发计划 V1.0", f"@{md_path}"],
                    capture_output=True, text=True, timeout=30,
                )
                if cr.returncode == 0:
                    m = re.search(r"https://[^\s\"']+feishu[^\s\"']+", cr.stdout)
                    wiki_url = m.group(0) if m else "(已创建)"
                else:
                    logger.warning("Wiki 创建失败: %s", (cr.stdout or cr.stderr or "")[:200])
            except Exception as exc:
                logger.warning("Wiki 创建异常: %s", exc)

        # ── 构建返回摘要 ───────────────────────────────────────────
        ms_lines = []
        for m in milestones:
            abs_date = (datetime.strptime(t0_date, "%Y-%m-%d") + timedelta(weeks=m["week"])).strftime("%Y-%m-%d")
            tag = "T0" if m["week"] == 0 else f"T+{m['week']}W"
            ms_lines.append(f"  {tag} ({abs_date})  {m['name']}")
        msg = (
            f"项目计划甘特图已生成！\n\n"
            f"项目：{project_name}\n"
            f"T0：{t0_date}  |  总周期：{total_weeks}W\n\n"
            f"里程碑：\n" + "\n".join(ms_lines) + "\n\n"
            f"本地 HTML：{html_path}"
        )
        if png_path:
            msg += f"\n截图 PNG：{png_path}"
        if html_url:
            msg += f"\n飞书 HTML：{html_url}"
        if png_url:
            msg += f"\n飞书 PNG：{png_url}"
        if wiki_url:
            msg += f"\n飞书 Wiki：{wiki_url}"
        return msg

    # ── 飞书项目 MCP ───────────────────────────────────────────────

    def _feishu_project(
        self,
        action: str,
        project_key: str = "",
        work_item_id: str = "",
        work_item_type: str = "",
        fields_list: list = None,
        fields_dict: dict = None,
        fields_updates: list = None,
        view_id: str = "",
        page_num: int = 1,
        nodes: list = None,
        node_id: str = "",
        moql: str = "",
        user_keys: list = None,
        start_date: str = "",
        end_date: str = "",
    ) -> str:
        try:
            mcp = FeishuProjectMCP()
        except ValueError as e:
            return f"飞书项目 MCP 未配置：{e}"

        try:
            if action == "get_workitem_brief":
                if not project_key or not work_item_id:
                    return "错误：get_workitem_brief 需要 project_key 和 work_item_id"
                return mcp.get_workitem_brief(project_key, work_item_id, fields=fields_list)

            if action == "get_workitem_info":
                if not project_key or not work_item_type:
                    return "错误：get_workitem_info 需要 project_key 和 work_item_type"
                return mcp.get_workitem_info(project_key, work_item_type)

            if action == "get_view_detail":
                if not project_key or not view_id:
                    return "错误：get_view_detail 需要 project_key 和 view_id"
                return mcp.get_view_detail(project_key, view_id,
                                           fields=fields_list, page_num=page_num)

            if action == "create_workitem":
                if not project_key or not work_item_type or not fields_dict:
                    return "错误：create_workitem 需要 project_key、work_item_type 和 fields_dict"
                return mcp.create_workitem(project_key, work_item_type, fields_dict)

            if action == "update_field":
                if not project_key or not work_item_id or not fields_updates:
                    return "错误：update_field 需要 project_key、work_item_id 和 fields_updates"
                return mcp.update_field(project_key, work_item_id, fields_updates)

            if action == "finish_node":
                if not project_key or not work_item_id or not nodes:
                    return "错误：finish_node 需要 project_key、work_item_id 和 nodes"
                return mcp.finish_node(project_key, work_item_id, nodes)

            if action == "get_node_detail":
                if not project_key or not work_item_id or not node_id:
                    return "错误：get_node_detail 需要 project_key、work_item_id 和 node_id"
                return mcp.get_node_detail(project_key, work_item_id, node_id)

            if action == "search_by_mql":
                if not moql:
                    return "错误：search_by_mql 需要 moql 参数"
                return mcp.search_by_mql(moql)

            if action == "list_schedule":
                if not project_key or not user_keys or not start_date or not end_date:
                    return "错误：list_schedule 需要 project_key、user_keys、start_date 和 end_date"
                return mcp.list_schedule(project_key, user_keys, start_date, end_date)

            return f"未知 feishu_project action: {action}"
        except Exception as e:
            logger.error("feishu_project[%s] 出错: %s", action, e, exc_info=True)
            return f"飞书项目操作失败（{action}）：{e}"

    # ── GB 国标检索 ────────────────────────────────────────────────

    def _search_gb_standard(self, query: str, standard: str) -> str:
        """在本地 GB 国标文档中全文搜索，返回匹配条款及上下文"""
        _DOCS = {
            "l2pp": os.path.join(
                _SKILL_BASE_DIR, "l2pp-gb-helper", "references",
                "智能网联汽车_组合驾驶辅助系统安全要求.md",
            ),
            "ads_l3l4": os.path.join(
                _SKILL_BASE_DIR, "ads-l3-l4-gb", "references", "GB_ADS.md"
            ),
        }
        _NAMES = {
            "l2pp": "《智能网联汽车 组合驾驶辅助系统安全要求》（L2++ GB报批稿）",
            "ads_l3l4": "《智能网联汽车 自动驾驶系统安全要求》（L3/L4 征求意见稿）",
        }
        doc_path = _DOCS.get(standard)
        doc_name = _NAMES.get(standard, standard)

        if not doc_path:
            return f"未知标准类型: {standard}，可选值：l2pp / ads_l3l4"
        if not os.path.exists(doc_path):
            return (
                f"[错误] {doc_name} 文档未找到。\n"
                f"预期路径：{doc_path}\n"
                "请确认 skill 已完整安装，或将文档放置到上述路径。"
            )

        try:
            result = subprocess.run(
                ["grep", "-n", "-i", "-C", "8", query, doc_path],
                capture_output=True, text=True, timeout=30,
            )
        except Exception as e:
            return f"[错误] 搜索失败：{e}"

        output = result.stdout.strip()
        if not output:
            return (
                f"在 {doc_name} 中未找到「{query}」相关内容。\n"
                "建议换用更短的核心词，如 '接管'、'TTC'、'4.6' 等。"
            )

        lines = output.split("\n")
        if len(lines) > 150:
            lines = lines[:150]
            lines.append("\n[...结果过多，已截断。请使用更精确的关键词缩小范围...]")

        return f"## {doc_name} 搜索结果（关键词：{query}）\n\n" + "\n".join(lines)

    # ── AI 多角色辩论 ──────────────────────────────────────────────

    def _debate(
        self,
        topic: str,
        context: str = "",
        persona_a: str = "",
        persona_b: str = "",
        rounds: int = 2,
        output_chat_id: str = "",
    ) -> str:
        """两个 AI 角色交替辩论，第三个角色总结结论。"""
        import anthropic as _anthropic
        import os as _os

        rounds = max(1, min(3, int(rounds)))

        sys_a = (
            f"你是{persona_a}。" if persona_a
            else "你是一位资深分析师，善于从多维度深入分析问题，给出全面、有深度的见解。"
        )
        sys_b = (
            f"你是{persona_b}。" if persona_b
            else "你是一位批判性思维专家，擅长发现方案漏洞、指出被忽视的风险，并提出建设性的改进意见。"
        )
        tone = "\n请用中文回答，控制在 300 字以内，观点清晰有重点，不要泛泛而谈。"

        _client = _anthropic.Anthropic(api_key=_os.environ.get("ANTHROPIC_API_KEY", ""))
        _model = "claude-sonnet-4-6"

        def _call(system: str, messages: list) -> str:
            resp = _client.messages.create(
                model=_model,
                max_tokens=1024,
                system=system + tone,
                messages=messages,
            )
            return "".join(b.text for b in resp.content if hasattr(b, "text")).strip()

        # 构建初始问题
        initial = f"话题：{topic}"
        if context:
            initial += f"\n\n背景信息：\n{context}"

        transcript: list[tuple[str, str]] = []  # (role_label, content)
        label_a = persona_a.split("，")[0] if persona_a else "分析师"
        label_b = persona_b.split("，")[0] if persona_b else "审查者"

        # Agent A 的对话历史
        msgs_a: list[dict] = [{"role": "user", "content": initial}]
        # Agent B 的对话历史（每轮重建）
        msgs_b: list[dict] = []

        for r in range(rounds):
            # ── Agent A 发言 ──
            text_a = _call(sys_a, msgs_a)
            transcript.append((label_a, text_a))
            msgs_a.append({"role": "assistant", "content": text_a})

            # ── Agent B 发言 ──
            if not msgs_b:
                msgs_b = [{"role": "user", "content": f"{initial}\n\n{label_a}的分析：\n{text_a}"}]
            else:
                msgs_b.append({"role": "user", "content": f"{label_a}的最新回应：\n{text_a}"})
            text_b = _call(sys_b, msgs_b)
            transcript.append((label_b, text_b))
            msgs_b.append({"role": "assistant", "content": text_b})

            # A 下一轮看到 B 的质疑
            if r < rounds - 1:
                msgs_a.append({
                    "role": "user",
                    "content": f"{label_b}的质疑：\n{text_b}\n\n请回应上述观点并进一步完善你的分析。",
                })

        # ── 仲裁者总结 ──
        debate_text = "\n\n".join(f"【{role}】\n{content}" for role, content in transcript)
        arbiter_prompt = (
            f"以下是关于「{topic}」的讨论记录：\n\n{debate_text}\n\n"
            "请：\n"
            "1. 用 1-2 句话概括双方的核心分歧\n"
            "2. 列出 3-5 条具体可执行的结论/建议（编号列表）\n"
            "3. 指出最需要关注的 1-2 个风险点\n"
            "格式清晰，控制在 400 字以内，用中文。"
        )
        conclusion = _call(
            "你是一位专业的讨论主持人，擅长总结多方观点、提炼共识、给出可执行建议。",
            [{"role": "user", "content": arbiter_prompt}],
        )

        # ── 格式化报告 ──
        lines = [f"**AI 多角色讨论：{topic}**\n"]
        for i, (role, content) in enumerate(transcript):
            round_num = i // 2 + 1
            marker = "▶" if role == label_a else "◀"
            lines.append(f"{marker} 第{round_num}轮 [{role}]\n{content}\n")
        lines.append(f"📋 综合结论\n{conclusion}")
        report = "\n".join(lines)

        if output_chat_id:
            try:
                self.feishu.send_text_to_chat(output_chat_id, report)
                return report + f"\n\n[已发送至群聊]"
            except Exception as e:
                return report + f"\n\n[发送群聊失败：{e}]"
        return report

    # ── Skill 搜索 ─────────────────────────────────────────────────

    def _find_skills(self, query: str, source: str = "both") -> str:
        """搜索公共 + 内部 skill 库，返回匹配结果供 AI 推荐"""
        import urllib.request
        import urllib.parse
        import re as _re
        from concurrent.futures import ThreadPoolExecutor, as_completed

        def _search_public() -> str:
            try:
                result = subprocess.run(
                    ["npx", "--yes", "skills", "find", query],
                    capture_output=True, text=True, timeout=25,
                    env={**os.environ, "NO_COLOR": "1"},
                )
                out = _re.sub(r"\x1b\[[0-9;]*m", "", (result.stdout or "").strip())
                if out:
                    lines = out.split("\n")
                    if len(lines) > 60:
                        lines = lines[:60] + ["[...已截断...]"]
                    return "## 公共技能库（skills.sh）搜索结果\n\n" + "\n".join(lines)
                return f"## 公共技能库（skills.sh）\n\n未找到与「{query}」相关的 skill。"
            except Exception as e:
                return f"## 公共技能库（skills.sh）\n\n搜索失败：{e}"

        def _search_internal() -> str:
            try:
                all_skills: list[dict] = []
                page = 1
                # 先用关键词做服务端搜索（快速路径）
                while True:
                    url = ("https://mcp-hub.momenta.works/skills/api/search?"
                           + urllib.parse.urlencode({"q": query, "page": page, "page_size": 100}))
                    req = urllib.request.Request(url, headers={"Accept": "application/json"})
                    with urllib.request.urlopen(req, timeout=10) as resp:
                        data = json.loads(resp.read())
                    skills = data.get("skills", [])
                    all_skills.extend(skills)
                    if len(all_skills) >= data.get("total", 0) or not skills:
                        break
                    page += 1

                # 结果太少时补全量列表（让 AI 自己判断）
                if len(all_skills) < 5:
                    all_skills = []
                    page = 1
                    while True:
                        url = ("https://mcp-hub.momenta.works/skills/api/search?"
                               + urllib.parse.urlencode({"q": "", "page": page, "page_size": 100}))
                        req = urllib.request.Request(url, headers={"Accept": "application/json"})
                        with urllib.request.urlopen(req, timeout=10) as resp:
                            data = json.loads(resp.read())
                        skills = data.get("skills", [])
                        all_skills.extend(skills)
                        if len(all_skills) >= data.get("total", 0) or not skills:
                            break
                        page += 1
                    header = f"## 内部技能库（mcp-hub.momenta.works）全量列表（共 {len(all_skills)} 个）\n请根据用户需求「{query}」从以下列表中识别最相关的 skill：\n"
                else:
                    header = f"## 内部技能库（mcp-hub.momenta.works）搜索结果（共 {len(all_skills)} 个匹配）\n"

                if not all_skills:
                    return "## 内部技能库\n\n未能获取技能列表。"

                lines = [header]
                for sk in all_skills:
                    name = sk.get("name", "")
                    desc = sk.get("description", "").replace("\n", " ")[:120]
                    source_url = sk.get("source", "")
                    likes = sk.get("likes_count", 0)
                    lines.append(f"- **{name}** ⭐{likes}｜{desc}")
                    if source_url:
                        lines.append(f"  安装：`npx skills install {source_url} --skill {name}`")
                return "\n".join(lines)
            except Exception as e:
                return f"## 内部技能库\n\n获取失败：{e}"

        # 并行执行两个搜索
        tasks: dict[str, str] = {}
        with ThreadPoolExecutor(max_workers=2) as pool:
            futures = {}
            if source in ("both", "public"):
                futures["public"] = pool.submit(_search_public)
            if source in ("both", "internal"):
                futures["internal"] = pool.submit(_search_internal)
            for key, fut in futures.items():
                tasks[key] = fut.result()

        sections = []
        if "public" in tasks:
            sections.append(tasks["public"])
        if "internal" in tasks:
            sections.append(tasks["internal"])
        return "\n\n---\n\n".join(sections)

    # ── NP&AEB 制动计算 ────────────────────────────────────────────

    def _np_aeb_calc(
        self,
        mode: str,
        v0: float,
        limits: list,
        v_aeb: float | None = None,
        s_det: float | None = None,
        t_delay: float = 0.5,
        aeb_a: float = 9.0,
        aeb_j: float = 20.0,
    ) -> str:
        """调用 calc_braking.py 进行 NP&AEB 制动模型精确计算"""
        _SCRIPT = os.path.join(
            os.path.expanduser("~"), "项目管理", "AI Agent",
            ".agents", "skills", "np-aeb-braking-calc", "scripts", "calc_braking.py"
        )
        if not os.path.exists(_SCRIPT):
            return "错误：calc_braking.py 未找到，请确认 np&aeb-braking-calc skill 已安装。"
        cmd = [
            sys.executable, _SCRIPT,
            "--mode", mode,
            "--v0", str(v0),
            "--t_delay", str(t_delay),
            "--aeb_a", str(aeb_a),
            "--aeb_j", str(aeb_j),
            "--limits", json.dumps(limits, ensure_ascii=False),
        ]
        if mode == "distance":
            if v_aeb is None:
                return "错误：distance 模式需要提供 v_aeb 参数。"
            cmd += ["--v_aeb", str(v_aeb)]
        else:
            if s_det is None:
                return "错误：latest_aeb 模式需要提供 s_det 参数。"
            cmd += ["--s_det", str(s_det)]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            if result.returncode != 0:
                return f"计算失败：{result.stderr.strip() or result.stdout.strip()}"
            return result.stdout.strip()
        except Exception as e:
            return f"计算异常：{e}"

    # ── 车辆 FMP 查询 ───────────────────────────────────────────────

    def _check_fmp_vehicles(
        self,
        project: str = "Project-MXS",
        hours: float = 8,
    ) -> str:
        """查询 FMP 平台指定项目的空闲/占用车辆。"""
        try:
            from feishu.fmp import check_session_valid, query_idle_vehicles
        except ImportError as e:
            return f"FMP 模块导入失败：{e}"

        if not check_session_valid():
            return (
                "⚠️ FMP 尚未授权，请先运行登录脚本：\n"
                "```\n"
                "python C:/Users/ryan.li/personal-assistant/scripts/fmp_login.py\n"
                "```\n"
                "在弹出的浏览器中用 Momenta 账号登录 FMP，登录完成后 bot 即可自动查询。"
            )

        result = query_idle_vehicles(project=project, end_hours=int(hours))

        if result.get("error"):
            return f"FMP 查询失败：{result['error']}"

        idle = result.get("idle", [])
        busy = result.get("busy", [])
        total = result.get("total", len(idle) + len(busy))
        q_time = result.get("query_time", "")

        lines = [f"# FMP 车辆状态 — {project}"]
        lines.append(f"查询时间段：{q_time}，共 {total} 辆")
        lines.append("")

        if idle:
            lines.append(f"## ✅ 空闲车辆（{len(idle)} 辆）")
            lines.append("| 车辆 ID | 车牌 | 车型 | 维保状态 |")
            lines.append("|--------|------|------|----------|")
            for c in idle:
                lines.append(
                    f"| {c['car_id']} | {c['car_plate']} | {c['car_type']} | {c['maintain_status'] or '正常'} |"
                )
        else:
            lines.append("## ✅ 空闲车辆：暂无")

        lines.append("")
        if busy:
            lines.append(f"## 🔴 占用中（{len(busy)} 辆）")
            for c in busy:
                lines.append(f"- {c['car_plate']} ({c['car_type']})")

        return "\n".join(lines)

    # ── 故障码查询 ───────────────────────────────────────────────────

    _BAG_FAULT_SCRIPT = "/mnt/c/Users/ryan.li/Desktop/bag-fault-query/bag_fault_query.py"
    _BAG_FAULT_PYTHON = "/root/bagenv/bin/python3"

    def _query_bag_fault(self, vin: str, time_points: str, size: int = 20) -> str:
        """通过 WSL 运行 bag_fault_query.py 查询 /mff_md/enable_signal_cmd 故障。"""
        vin = vin.strip().upper()
        if len(vin) != 17:
            return f"[错误] VIN 必须是 17 位，当前: {vin!r}"

        time_args = time_points.strip().split()
        cmd = [
            "wsl", "-d", "Ubuntu", "--",
            "bash", "-c",
            f"{self._BAG_FAULT_PYTHON} {self._BAG_FAULT_SCRIPT} "
            f"{vin} {' '.join(time_args)} --size {size} 2>&1 | grep -v pybagmining",
        ]
        logger.info("_query_bag_fault: %s %s size=%d", vin, time_args, size)
        try:
            result = subprocess.run(cmd, capture_output=True, timeout=180)
            output = result.stdout.decode("utf-8", errors="replace").strip()
            stderr  = result.stderr.decode("utf-8", errors="replace").strip()
            if not output and stderr:
                return f"[错误] 查询失败：\n{stderr}"
            if result.returncode != 0 and not output:
                return f"[错误] 脚本退出码 {result.returncode}，无输出"
            return output or "查询完成，无输出"
        except subprocess.TimeoutExpired:
            return "[错误] 查询超时（>180s），ESS 或 bag 服务可能繁忙，请稍后重试"
        except FileNotFoundError:
            return "[错误] 未找到 wsl 命令，请确认 WSL Ubuntu 已安装"
        except Exception as e:
            return f"[错误] 查询异常：{e}"

    # ── 车辆预约 ────────────────────────────────────────────────────

    _FLEET_BOT_OPEN_ID = "ou_d9fc64ef0fd2949a226b50d79ea0ee09"
    _SAFETY_GROUP_CHAT_ID = "oc_c75f4dca2c72189e94f84dc555813939"

    def _book_vehicle(
        self,
        vehicle_id: str,
        time_range: str,
        task_name: str = "集成测试",
        project: str = "",
    ) -> str:
        """私发 Fleet-Bot 发送车辆预约请求（以 Ryan 用户身份发送）。"""
        body_parts = [f"约车并审批{vehicle_id}，{task_name}，{time_range}"]
        if project:
            body_parts.append(f"归属{project}")
        msg = "，".join(body_parts)
        logger.info("_book_vehicle: 私发 Fleet-Bot: %s", msg)
        try:
            from feishu_sync.api import send_message as _feishu_sync_send
            import json as _json
            result = _feishu_sync_send(
                receive_id=self._FLEET_BOT_OPEN_ID,
                msg_type="text",
                content=_json.dumps({"text": msg}),
                receive_id_type="open_id",
                as_user=True,
            )
            mid = result.get("message_id", "") if isinstance(result, dict) else ""
        except Exception as e:
            logger.error("_book_vehicle 私发失败: %s", e)
            mid = ""
        if mid:
            return (
                f"已私发 Fleet-Bot 发送车辆预约请求，Fleet-Bot 将自动处理预约和审批。\n"
                f"- 车辆：{vehicle_id}\n"
                f"- 任务：{task_name}\n"
                f"- 时间：{time_range}\n"
                + (f"- 归属：{project}\n" if project else "")
                + f"message_id={mid}"
            )
        return "预约消息发送失败，请检查 feishu-sync 是否已授权。"


# ------------------------------------------------------------------ #
# PPT 生成脚本构建
# ------------------------------------------------------------------ #

def _build_ppt_script(data: dict) -> str:
    """生成 Momenta 风格 PPT 的 Python 脚本"""
    data_json = json.dumps(data, ensure_ascii=False)
    return textwrap.dedent(f"""\
        # -*- coding: utf-8 -*-
        import sys, os, json, platform, subprocess

        # 自动安装依赖
        for pkg in ["pptx", "lxml"]:
            try:
                __import__(pkg)
            except ImportError:
                subprocess.check_call([sys.executable, "-m", "pip", "install",
                    "python-pptx" if pkg == "pptx" else pkg, "-q"])

        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import lxml.etree as etree
        from pptx.oxml.ns import qn

        DATA = json.loads({data_json!r})

        # ── 颜色常量 ────────────────────────────────────────
        BLUE   = RGBColor(0x00, 0x68, 0xE9)
        ORANGE = RGBColor(0xED, 0x7D, 0x31)
        DARK   = RGBColor(0x1F, 0x23, 0x29)
        GRAY   = RGBColor(0x75, 0x78, 0x7E)
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        BG     = RGBColor(0xF5, 0xF7, 0xFA)

        # 幻灯片尺寸（16:9 宽屏）
        SLIDE_W = Inches(13.33)
        SLIDE_H = Inches(7.5)

        CJK_FONT = "PingFang SC" if platform.system() == "Darwin" else "微软雅黑"

        def set_font(run, face):
            rPr = run._r.get_or_add_rPr()
            for tag in (qn("a:latin"), qn("a:ea")):
                el = rPr.find(tag)
                if el is None:
                    el = etree.SubElement(rPr, tag)
                el.set("typeface", face)

        def add_text_box(slide, text, left, top, width, height,
                         size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                         word_wrap=True):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            txBox.word_wrap = word_wrap
            tf = txBox.text_frame
            tf.word_wrap = word_wrap
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            set_font(run, CJK_FONT)
            return txBox

        def add_rect(slide, left, top, width, height, fill_color):
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.fill.background()
            return shape

        template_path = DATA.get("template_path", "")
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            # 移除模版中的所有幻灯片，保留母版/主题/布局
            # 必须同时删除 relationship，否则孤立的 slide part 会留在包里
            # 导致 PowerPoint 打开时报"格式损坏"
            _R_ID_KEY = '{{http://schemas.openxmlformats.org/officeDocument/2006/relationships}}id'
            sldIdLst = prs.slides._sldIdLst
            for sldId in list(sldIdLst):
                rId = sldId.get(_R_ID_KEY)
                sldIdLst.remove(sldId)
                if rId:
                    prs.part.rels._rels.pop(rId, None)
        else:
            prs = Presentation()
            prs.slide_width  = SLIDE_W
            prs.slide_height = SLIDE_H

        blank_layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]  # blank

        # ── 封面 ────────────────────────────────────────────
        cover = prs.slides.add_slide(blank_layout)

        # 深蓝背景
        add_rect(cover, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0x00, 0x1A, 0x57))

        # 右侧装饰条
        strip_w = Inches(0.5)
        add_rect(cover, SLIDE_W - strip_w, 0, strip_w, SLIDE_H, BLUE)

        # 主标题
        add_text_box(
            cover, DATA["title"],
            Inches(1), Inches(2.2), Inches(10), Inches(1.5),
            size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
        )

        # 副标题
        if DATA.get("subtitle"):
            add_text_box(
                cover, DATA["subtitle"],
                Inches(1), Inches(3.9), Inches(10), Inches(0.8),
                size=20, bold=False, color=RGBColor(0xB0, 0xC4, 0xFF),
            )

        # 日期
        from datetime import date
        add_text_box(
            cover, date.today().strftime("%Y.%m"),
            Inches(1), Inches(6.2), Inches(4), Inches(0.5),
            size=14, color=GRAY,
        )

        # ── 内容页 ───────────────────────────────────────────
        for slide_data in DATA["slides"]:
            sl = prs.slides.add_slide(blank_layout)

            # 白色背景
            add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, WHITE)

            # 顶部蓝色标题栏
            header_h = Inches(1.0)
            add_rect(sl, 0, 0, SLIDE_W, header_h, BLUE)

            # 页面标题
            add_text_box(
                sl, slide_data["title"],
                Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                size=24, bold=True, color=WHITE,
            )

            # 要点列表
            bullets = slide_data.get("bullets", [])
            y = Inches(1.2)
            line_h = Inches(0.45)
            for i, bullet in enumerate(bullets):
                # 圆点
                add_rect(sl,
                    Inches(0.5), y + Inches(0.12),
                    Inches(0.08), Inches(0.08),
                    ORANGE,
                )
                add_text_box(
                    sl, bullet,
                    Inches(0.7), y, Inches(11.8), line_h,
                    size=14, color=DARK,
                )
                y += line_h
                if y > Inches(6.8):  # 防止超出页面
                    break

            # 备注
            note = slide_data.get("note", "")
            if note:
                add_text_box(
                    sl, note,
                    Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
                    size=10, color=GRAY,
                )

        prs.save(DATA["output_path"])
        print(f"OK: {{DATA['output_path']}}")
    """)


# ------------------------------------------------------------------ #
# 辅助函数
# ------------------------------------------------------------------ #

def _convert_timestamps(obj):
    """递归将 {sec, nsec} 结构转换为浮点秒，便于阅读"""
    if isinstance(obj, dict):
        if set(obj.keys()) == {"sec", "nsec"}:
            return round(obj["sec"] + obj["nsec"] / 1e9, 6)
        return {k: _convert_timestamps(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_convert_timestamps(i) for i in obj]
    return obj


def _parse_dt(dt_str: str) -> datetime:
    for fmt in ("%Y-%m-%dT%H:%M:%S", "%Y-%m-%dT%H:%M", "%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M"):
        try:
            naive = datetime.strptime(dt_str.split("+")[0].split("Z")[0], fmt)
            return TZ_SHANGHAI.localize(naive)
        except ValueError:
            continue
    raise ValueError(f"无法解析时间: {dt_str}")
