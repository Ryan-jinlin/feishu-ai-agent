"""Feishu Open API 封装（认证、消息、日历、用户搜索）"""
from __future__ import annotations

import os
import re
import time
import json
import logging
from datetime import datetime
import pytz
import requests

logger = logging.getLogger(__name__)

FEISHU_HOST = "https://open.feishu.cn"
TZ_SHANGHAI = pytz.timezone("Asia/Shanghai")

# 持久化应用日历 ID，避免每次重启都创建新日历
_CALENDAR_ID_FILE = os.path.join(os.path.dirname(__file__), "..", ".feishu_calendar_id")

# 用户 IM OAuth token（由 scripts/authorize_user_im.py 生成，用于读取 Bot 不在的群）
_USER_IM_TOKEN_FILE = os.path.join(os.path.dirname(__file__), "..", ".user_im_token.json")

# Momenta 蓝：飞书卡片 template 使用 "blue"（最接近 #0066FF）
_CARD_TEMPLATE = "blue"


_LARK_MD_MAX = 2000  # 每个 lark_md element 的安全内容上限（字符）


def _split_lark_md(content: str) -> list[str]:
    """将过长的 lark_md 内容按段落切割，每块不超过 _LARK_MD_MAX 字符。"""
    if len(content) <= _LARK_MD_MAX:
        return [content]
    chunks: list[str] = []
    current = ""
    for para in content.split('\n\n'):
        candidate = (current + '\n\n' + para).lstrip('\n') if current else para
        if len(candidate) > _LARK_MD_MAX:
            if current:
                chunks.append(current)
            # 单段超长则按行再拆
            current = para[:_LARK_MD_MAX]
        else:
            current = candidate
    if current:
        chunks.append(current)
    return chunks or [content[:_LARK_MD_MAX]]


def _append_md(elements: list[dict], content: str) -> None:
    """将 lark_md 内容追加为一个或多个 div 元素（自动切割超长块）。"""
    content = content.strip()
    if not content:
        return
    for chunk in _split_lark_md(content):
        chunk = chunk.strip()
        if chunk:
            elements.append({"tag": "div", "text": {"tag": "lark_md", "content": chunk}})


def _section_header(heading: str) -> dict:
    """生成 H2 章节标题元素：灰底 column_set，视觉效果接近飞书文档的分节标题。"""
    return {
        "tag": "column_set",
        "flex_mode": "none",
        "background_style": "grey",
        "columns": [{
            "tag": "column",
            "width": "weighted",
            "weight": 1,
            "elements": [{
                "tag": "div",
                "text": {"tag": "lark_md", "content": f"**{heading}**"},
            }],
        }],
    }


def _md_to_feishu_card(text: str) -> dict:
    """将 Markdown 文本转换为飞书 Interactive Card JSON。

    转换规则：
    - ``# Title``   → 卡片 header 蓝色标题栏（否则默认 "AI 助理"）
    - ``## Heading`` → 灰底 column_set 分节标题（加粗，无 ## 前缀）
    - ``### / ####`` → 加粗内联（无 # 前缀）
    - ``---``        → hr 分割线
    - lark_md element 超过 2000 字符时自动切割（防 Feishu API 报错）
    """
    # 1. 提取 H1 作为卡片 header 标题
    title = "AI 助理"
    h1_match = re.match(r'^#\s+(.+?)$', text, re.MULTILINE)
    if h1_match:
        title = h1_match.group(1).strip()
        text = text[:h1_match.start()] + text[h1_match.end():]
        text = text.lstrip('\n')

    elements: list[dict] = []
    pending: list[str] = []  # 待积累的普通行

    def flush_pending() -> None:
        _append_md(elements, '\n'.join(pending))
        pending.clear()

    for line in text.split('\n'):
        m2 = re.match(r'^##\s+(.+?)$', line)
        m3p = re.match(r'^#{3,6}\s+(.+?)$', line)
        is_hr = bool(re.match(r'^\s*-{3,}\s*$', line))

        if m2:
            flush_pending()
            elements.append(_section_header(m2.group(1).strip()))
        elif m3p:
            # H3~H6：加粗内联，融入当前内容块
            pending.append(f"**{m3p.group(1).strip()}**")
        elif is_hr:
            flush_pending()
            elements.append({"tag": "hr"})
        else:
            pending.append(line)

    flush_pending()

    if not elements:
        _append_md(elements, text.strip() or "完成")

    return {
        "config": {"wide_screen_mode": True},
        "header": {
            "title": {"tag": "plain_text", "content": title},
            "template": _CARD_TEMPLATE,
        },
        "elements": elements,
    }


class FeishuClient:
    def __init__(self, app_id: str, app_secret: str):
        self.app_id = app_id
        self.app_secret = app_secret
        self._token: str = ""
        self._token_expires_at: float = 0
        self._cached_calendar_id: str = ""
        # 用户目录缓存（后台线程预热，避免首次搜索慢）
        self._user_cache: list[dict] = []   # [{name, open_id, email}]
        self._user_cache_ts: float = 0       # 上次构建时间戳（0=未构建）
        self._user_cache_ready = False
        import threading as _threading
        _threading.Thread(target=self._warm_user_cache, daemon=True).start()

    # ------------------------------------------------------------------ #
    # 认证
    # ------------------------------------------------------------------ #

    def _get_token(self) -> str:
        if time.time() < self._token_expires_at - 60:
            return self._token
        resp = requests.post(
            f"{FEISHU_HOST}/open-apis/auth/v3/tenant_access_token/internal",
            json={"app_id": self.app_id, "app_secret": self.app_secret},
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            raise RuntimeError(f"获取 Token 失败: {data}")
        self._token = data["tenant_access_token"]
        self._token_expires_at = time.time() + data.get("expire", 7200)
        return self._token

    def _headers(self) -> dict:
        return {"Authorization": f"Bearer {self._get_token()}", "Content-Type": "application/json"}

    # ------------------------------------------------------------------ #
    # Bot 自身信息
    # ------------------------------------------------------------------ #

    def get_bot_open_id(self) -> str:
        """获取当前应用 Bot 的 open_id（用于群聊 @mention 过滤）"""
        try:
            resp = requests.get(
                f"{FEISHU_HOST}/open-apis/bot/v3/info",
                headers=self._headers(),
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
            open_id = data.get("bot", {}).get("open_id", "")
            if open_id:
                logger.info("Bot open_id: %s", open_id)
            return open_id
        except Exception as e:
            logger.warning("获取 Bot open_id 失败: %s", e)
            return ""

    # ------------------------------------------------------------------ #
    # 用户搜索
    # ------------------------------------------------------------------ #

    def search_users(self, keyword: str, page_size: int = 5) -> list[dict]:
        """按关键词搜索飞书用户，返回 [{name, open_id, email}]

        只使用 Contact API（tenant token）确保返回 Bot 应用正确的 open_id。
        注意：/search/v1/user（user_access_token）返回的 open_id 是 feishu-sync
        OAuth app 的 open_id，与 Bot app 不同（跨应用），不能用于 IM 操作。
        """
        # 1. 使用用户目录缓存（tenant token，open_id 正确）
        if self._user_cache_ready:
            # 检查是否需要刷新缓存
            if time.time() - self._user_cache_ts > self._USER_CACHE_TTL:
                import threading as _threading
                _threading.Thread(target=self._warm_user_cache, daemon=True).start()
            def _norm(s: str) -> str:
                """去掉全角/半角空格，转小写，便于模糊匹配"""
                return s.replace('\u3000', '').replace(' ', '').lower()
            kw = _norm(keyword)
            results = [u for u in self._user_cache if kw in _norm(u["name"]) or kw in _norm(u.get("en_name", ""))]
            return results[:page_size]

        # 2. 缓存未就绪，回退到实时遍历（速度较慢，仅启动后首次调用）
        return self._search_users_by_contact(keyword, page_size)

    _USER_CACHE_TTL = 7200  # 2小时刷新一次

    def _warm_user_cache(self) -> None:
        """后台线程：预热用户目录缓存（启动时自动调用）"""
        try:
            logger.info("开始预热用户目录缓存...")
            cache = self._fetch_all_users_via_contact()
            self._user_cache = cache
            self._user_cache_ts = time.time()
            self._user_cache_ready = True
            logger.info("用户目录缓存预热完成，共 %d 名用户", len(cache))
        except Exception as e:
            logger.warning("用户目录缓存预热失败: %s", e)

    def _fetch_all_users_via_contact(self) -> list[dict]:
        """用 Contact API + tenant token 拉取全组织用户（最多 3 层部门树）"""
        all_users: list[dict] = []
        seen: set[str] = set()

        def _fetch_dept(dept_id: str) -> None:
            page_token = None
            while True:
                params = {"user_id_type": "open_id", "department_id": dept_id, "page_size": 50}
                if page_token:
                    params["page_token"] = page_token
                try:
                    resp = requests.get(
                        f"{FEISHU_HOST}/open-apis/contact/v3/users",
                        headers=self._headers(), params=params, timeout=10,
                    )
                    data = resp.json()
                except Exception:
                    break
                if data.get("code") != 0:
                    break
                for u in data.get("data", {}).get("items", []):
                    oid = u.get("open_id", "")
                    if oid and oid not in seen:
                        seen.add(oid)
                        all_users.append({
                            "name": u.get("name", ""),
                            "en_name": u.get("en_name", ""),
                            "open_id": oid,
                            "email": u.get("email", ""),
                        })
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token")

        def _child_depts(parent_id: str) -> list[str]:
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/contact/v3/departments",
                    headers=self._headers(),
                    params={"user_id_type": "open_id", "department_id_type": "open_department_id",
                            "parent_department_id": parent_id, "page_size": 50},
                    timeout=10,
                )
                return [d["open_department_id"] for d in resp.json().get("data", {}).get("items", [])]
            except Exception:
                return []

        for d1 in _child_depts("0"):
            _fetch_dept(d1)
            for d2 in _child_depts(d1):
                _fetch_dept(d2)
                for d3 in _child_depts(d2):
                    _fetch_dept(d3)

        return all_users

    def _search_users_by_contact(self, keyword: str, page_size: int) -> list[dict]:
        """用 Contact API（tenant token）遍历部门树，按名字关键词过滤用户"""
        def _norm(s: str) -> str:
            return s.replace('\u3000', '').replace(' ', '').lower()
        kw = _norm(keyword)
        found: list[dict] = []
        seen_open_ids: set[str] = set()

        def _fetch_users_in_dept(dept_id: str) -> None:
            """拉取某部门下的所有用户（非递归子部门）"""
            page_token = None
            while True:
                params = {"user_id_type": "open_id", "department_id": dept_id, "page_size": 50}
                if page_token:
                    params["page_token"] = page_token
                try:
                    resp = requests.get(
                        f"{FEISHU_HOST}/open-apis/contact/v3/users",
                        headers=self._headers(),
                        params=params,
                        timeout=10,
                    )
                    data = resp.json()
                except Exception:
                    break
                if data.get("code") != 0:
                    break
                for u in data.get("data", {}).get("items", []):
                    name = u.get("name", "")
                    en_name = u.get("en_name", "")
                    oid = u.get("open_id", "")
                    if oid in seen_open_ids:
                        continue
                    if kw in _norm(name) or kw in _norm(en_name):
                        seen_open_ids.add(oid)
                        found.append({
                            "name": name,
                            "en_name": en_name,
                            "open_id": oid,
                            "email": u.get("email", ""),
                            "department": "",
                        })
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token")

        def _get_child_depts(parent_id: str) -> list[str]:
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/contact/v3/departments",
                    headers=self._headers(),
                    params={
                        "user_id_type": "open_id",
                        "department_id_type": "open_department_id",
                        "parent_department_id": parent_id,
                        "page_size": 50,
                    },
                    timeout=10,
                )
                data = resp.json()
                return [d["open_department_id"] for d in data.get("data", {}).get("items", [])]
            except Exception:
                return []

        # BFS：遍历最多 3 层部门（根 → 一级 → 二级 → 三级）
        level1 = _get_child_depts("0")
        for d1 in level1:
            _fetch_users_in_dept(d1)
            if len(found) >= page_size:
                return found[:page_size]
            level2 = _get_child_depts(d1)
            for d2 in level2:
                _fetch_users_in_dept(d2)
                if len(found) >= page_size:
                    return found[:page_size]
                for d3 in _get_child_depts(d2):
                    _fetch_users_in_dept(d3)
                    if len(found) >= page_size:
                        return found[:page_size]

        return found[:page_size]

    def get_user_by_open_id(self, open_id: str) -> dict | None:
        """根据 open_id 获取用户信息"""
        resp = requests.get(
            f"{FEISHU_HOST}/open-apis/contact/v3/users/{open_id}",
            headers=self._headers(),
            params={"user_id_type": "open_id"},
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.warning("获取用户失败: %s", data)
            return None
        return data.get("data", {}).get("user")

    # ------------------------------------------------------------------ #
    # 消息
    # ------------------------------------------------------------------ #

    def reply_message(self, message_id: str, text: str) -> bool:
        """回复某条消息（纯文本）"""
        resp = requests.post(
            f"{FEISHU_HOST}/open-apis/im/v1/messages/{message_id}/reply",
            headers=self._headers(),
            json={"content": json.dumps({"text": text}), "msg_type": "text"},
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.error("回复消息失败: %s", data)
            return False
        return True

    def reply_card(self, message_id: str, text: str) -> bool:
        """以飞书消息卡片（Interactive Card）格式回复，将 Markdown 渲染为富文本。
        失败时自动降级为纯文本回复。
        """
        try:
            card = _md_to_feishu_card(text)
            card_json = json.dumps(card, ensure_ascii=False)
            logger.debug("reply_card elements=%d total_len=%d",
                         len(card.get("elements", [])), len(card_json))
            resp = requests.post(
                f"{FEISHU_HOST}/open-apis/im/v1/messages/{message_id}/reply",
                headers=self._headers(),
                json={"content": card_json, "msg_type": "interactive"},
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
            if data.get("code") != 0:
                logger.warning("回复卡片失败 code=%s msg=%s，降级为文本",
                               data.get("code"), data.get("msg", ""))
                return self.reply_message(message_id, text)
            return True
        except requests.HTTPError as exc:
            body = exc.response.text[:300] if exc.response else ""
            logger.warning("reply_card HTTP %s: %s，降级为文本",
                           exc.response.status_code if exc.response else "?", body)
            return self.reply_message(message_id, text)
        except Exception as exc:
            logger.warning("reply_card 异常，降级为文本: %s", exc)
            return self.reply_message(message_id, text)

    def download_message_resource(self, message_id: str, file_key: str) -> bytes | None:
        """下载消息中的图片资源，返回原始字节；失败返回 None。"""
        try:
            resp = requests.get(
                f"{FEISHU_HOST}/open-apis/im/v1/messages/{message_id}/resources/{file_key}",
                params={"type": "image"},
                headers=self._headers(),
                timeout=30,
            )
            if resp.status_code == 200:
                return resp.content
            logger.warning("download_message_resource 失败: %s %s", resp.status_code, resp.text[:200])
        except Exception as e:
            logger.warning("download_message_resource 异常: %s", e)
        return None

    def download_drive_file(self, obj_token: str) -> bytes | None:
        """下载飞书 Drive 文件，返回原始字节；失败返回 None。"""
        try:
            resp = requests.get(
                f"{FEISHU_HOST}/open-apis/drive/v1/files/{obj_token}/download",
                headers=self._headers(),
                stream=True,
                timeout=60,
            )
            if not resp.ok:
                logger.warning("download_drive_file 失败: %s %s", resp.status_code, resp.text[:200])
                return None
            return resp.content
        except Exception as e:
            logger.warning("download_drive_file 异常: %s", e)
            return None

    def get_drive_file_parent(self, obj_token: str) -> tuple[str, str]:
        """查询 Drive 文件所在的父目录。返回 (parent_node_token, parent_type)。
        parent_type 为 'explorer'（我的空间文件夹）或 'wiki'（知识库空间）。
        obj_token 为空时直接返回「我的空间」根目录。失败时返回 ('', 'explorer')。
        """
        # 有 obj_token 时，尝试通过 Drive file meta 接口查询父目录
        if obj_token:
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/drive/v1/files/{obj_token}",
                    headers=self._headers(),
                    timeout=10,
                )
                if resp.ok:
                    data = resp.json().get("data", {})
                    parent_token = data.get("parent_token", "")
                    parent_type  = data.get("parent_type",  "explorer")
                    if parent_token:
                        logger.info("get_drive_file_parent OK: %s -> %s (%s)", obj_token, parent_token, parent_type)
                        return parent_token, parent_type
            except Exception as e:
                logger.warning("get_drive_file_parent 异常: %s", e)

        # Fallback：获取「我的空间」根目录
        try:
            resp2 = requests.get(
                f"{FEISHU_HOST}/open-apis/drive/explorer/v2/root_folder/meta",
                headers=self._headers(),
                timeout=10,
            )
            if resp2.ok:
                root_token = resp2.json().get("data", {}).get("token", "")
                if root_token:
                    logger.info("get_drive_file_parent fallback 使用我的空间根目录: %s", root_token)
                    return root_token, "explorer"
        except Exception as e2:
            logger.warning("get_drive_file_parent fallback 异常: %s", e2)

        return "", "explorer"

    def upload_drive_file(
        self,
        filename: str,
        content: bytes,
        parent_token: str,
        parent_type: str = "explorer",
    ) -> tuple[str, str]:
        """上传文件到飞书 Drive 指定目录。返回 (file_token, 飞书链接)；失败返回 ('', '')。"""
        import io as _io

        # 上传时不能带 Content-Type: application/json，让 requests 自动设置 multipart boundary
        headers = {k: v for k, v in self._headers().items() if k.lower() != "content-type"}
        try:
            resp = requests.post(
                f"{FEISHU_HOST}/open-apis/drive/v1/files/upload_all",
                headers=headers,
                data={
                    "file_name":   filename,
                    "parent_type": parent_type,
                    "parent_node": parent_token,
                    "size":        str(len(content)),
                },
                files={"file": (filename, _io.BytesIO(content), "application/octet-stream")},
                timeout=120,
            )
            if resp.ok:
                file_token = resp.json().get("data", {}).get("file_token", "")
                url = f"https://momenta.feishu.cn/file/{file_token}" if file_token else ""
                logger.info("upload_drive_file 成功: %s -> %s", filename, url)
                return file_token, url
            else:
                logger.warning("upload_drive_file 失败: %s %s", resp.status_code, resp.text[:300])
                return "", ""
        except Exception as e:
            logger.warning("upload_drive_file 异常: %s", e)
            return "", ""

    @staticmethod
    def parse_pptx_text(data: bytes) -> str:
        """从 PPTX 字节数据中提取所有文字内容，按幻灯片分组。"""
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if slide_texts:
                parts.append(f"### 第{idx + 1}页\n" + "\n".join(slide_texts))
        return "\n\n".join(parts) or "(未提取到文字内容)"

    def send_text_to_user(self, open_id: str, text: str) -> str:
        """向用户发送私信（纯文本）。返回 message_id（成功）或空字符串（失败）。"""
        resp = requests.post(
            f"{FEISHU_HOST}/open-apis/im/v1/messages",
            headers=self._headers(),
            params={"receive_id_type": "open_id"},
            json={
                "receive_id": open_id,
                "content": json.dumps({"text": text}),
                "msg_type": "text",
            },
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.error("发送消息失败: %s", data)
            return ""
        return data.get("data", {}).get("message_id", "") or ""

    def send_text_to_chat(self, chat_id: str, text: str) -> str:
        """向群聊发送文本消息（text 中可含 <at user_id="open_id">姓名</at> @mention）。
        Bot 不在群内时自动 fallback 到用户 IM token（以用户身份发送）。
        返回 message_id（成功）或空字符串（失败）。"""
        payload = {
            "receive_id": chat_id,
            "content": json.dumps({"text": text}),
            "msg_type": "text",
        }
        params = {"receive_id_type": "chat_id"}
        # 先用 Bot token 尝试
        try:
            resp = requests.post(
                f"{FEISHU_HOST}/open-apis/im/v1/messages",
                headers=self._headers(),
                params=params,
                json=payload,
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
            if data.get("code") == 0:
                return data.get("data", {}).get("message_id", "") or ""
            logger.warning("send_text_to_chat Bot token 失败 code=%s，尝试用户 IM token", data.get("code"))
        except requests.HTTPError as e:
            if e.response is None or e.response.status_code not in (400, 403):
                raise
            logger.warning("send_text_to_chat Bot token HTTP %s，尝试用户 IM token",
                           e.response.status_code)
        # Fallback：用用户 IM token（Bot 不在群、但用户在群内）
        user_headers = self._user_im_headers()
        if not user_headers:
            logger.error("send_text_to_chat：Bot token 失败且无用户 IM token，无法发送")
            return ""
        resp2 = requests.post(
            f"{FEISHU_HOST}/open-apis/im/v1/messages",
            headers={**user_headers, "Content-Type": "application/json"},
            params=params,
            json=payload,
            timeout=10,
        )
        resp2.raise_for_status()
        data2 = resp2.json()
        if data2.get("code") != 0:
            logger.error("send_text_to_chat 用户 IM token 也失败: %s", data2)
            return ""
        logger.info("send_text_to_chat：已用用户 IM token 以用户身份发送到群 %s", chat_id)
        return data2.get("data", {}).get("message_id", "") or ""

    def recall_message(self, message_id: str) -> bool:
        """撤回（删除）指定消息 DELETE /open-apis/im/v1/messages/{message_id}"""
        resp = requests.delete(
            f"{FEISHU_HOST}/open-apis/im/v1/messages/{message_id}",
            headers=self._headers(),
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.error("撤回消息失败: %s", data)
            return False
        return True

    def send_card_to_user(self, open_id: str, card: dict) -> bool:
        """向用户发送消息卡片"""
        resp = requests.post(
            f"{FEISHU_HOST}/open-apis/im/v1/messages",
            headers=self._headers(),
            params={"receive_id_type": "open_id"},
            json={
                "receive_id": open_id,
                "content": json.dumps(card),
                "msg_type": "interactive",
            },
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.error("发送卡片失败: %s", data)
            return False
        return True

    # ------------------------------------------------------------------ #
    # 日历
    # ------------------------------------------------------------------ #

    def _get_or_create_app_calendar_id(self) -> str | None:
        """获取或创建应用自有的共享日历（使用 tenant_access_token，无需 user_access_token）"""
        if self._cached_calendar_id:
            return self._cached_calendar_id

        # 优先从持久化文件中读取，避免每次重启都创建新日历
        cal_id_file = os.path.abspath(_CALENDAR_ID_FILE)
        if os.path.exists(cal_id_file):
            try:
                with open(cal_id_file, encoding="utf-8") as f:
                    saved_id = f.read().strip()
                if saved_id:
                    self._cached_calendar_id = saved_id
                    logger.info("从文件加载应用日历 ID: %s", saved_id)
                    return saved_id
            except Exception as e:
                logger.warning("读取日历 ID 文件失败: %s", e)

        # 创建一个应用拥有的共享日历（tenant_access_token 可用）
        resp = requests.post(
            f"{FEISHU_HOST}/open-apis/calendar/v4/calendars",
            headers=self._headers(),
            json={"summary": "会议助理日历", "color": -1},
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.error("创建应用日历失败: %s", data)
            return None
        cal_id = data.get("data", {}).get("calendar", {}).get("calendar_id", "")
        if cal_id:
            self._cached_calendar_id = cal_id
            logger.info("应用日历已创建: %s", cal_id)
            try:
                with open(cal_id_file, "w", encoding="utf-8") as f:
                    f.write(cal_id)
            except Exception as e:
                logger.warning("保存日历 ID 文件失败: %s", e)
        return cal_id or None

    def create_calendar_event(
        self,
        title: str,
        description: str,
        start_dt: datetime,
        end_dt: datetime,
        location: str = "",
    ) -> dict | None:
        """
        创建日历事件（不含与会者，与会者通过 add_event_attendees 单独邀请）。
        返回 {"event_id": ..., "calendar_id": ...}，失败返回 None。
        """
        calendar_id = self._get_or_create_app_calendar_id()
        if not calendar_id:
            return None

        def to_feishu_time(dt: datetime) -> dict:
            if dt.tzinfo is None:
                dt = TZ_SHANGHAI.localize(dt)
            return {
                "timestamp": str(int(dt.timestamp())),
                "timezone": "Asia/Shanghai",
            }

        event_body: dict = {
            "summary": title,
            "description": description,
            "start_time": to_feishu_time(start_dt),
            "end_time": to_feishu_time(end_dt),
            "attendee_ability": "can_modify_event",  # 允许与会者编辑会议（邀请他人/修改时间和内容）
        }
        if location:
            event_body["location"] = {"name": location}

        resp = requests.post(
            f"{FEISHU_HOST}/open-apis/calendar/v4/calendars/{calendar_id}/events",
            headers=self._headers(),
            params={"user_id_type": "open_id"},
            json=event_body,
            timeout=15,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.error("创建日历事件失败: %s", data)
            return None

        event = data.get("data", {}).get("event", {})
        event_id = event.get("event_id", "")
        logger.info("日历事件已创建: %s (calendar: %s)", event_id, calendar_id)
        return {"event_id": event_id, "calendar_id": calendar_id}

    def add_event_attendees(
        self,
        calendar_id: str,
        event_id: str,
        attendee_open_ids: list[str],
        optional_open_ids: list[str] | None = None,
    ) -> bool:
        """通过 /attendees API 向日历事件添加与会者并发送正式邀请。
        分两次调用：必选与会者一次，可选（optional）与会者单独一次（失败则降级为普通成员）。
        """
        def _post_attendees(entries: list[dict]) -> bool:
            if not entries:
                return True
            try:
                resp = requests.post(
                    f"{FEISHU_HOST}/open-apis/calendar/v4/calendars/{calendar_id}/events/{event_id}/attendees",
                    headers=self._headers(),
                    params={"user_id_type": "open_id"},
                    json={"attendees": entries, "need_notification": True},
                    timeout=15,
                )
                resp.raise_for_status()
                data = resp.json()
                if data.get("code") != 0:
                    logger.error("添加与会者失败: %s", data)
                    return False
                return True
            except Exception as e:
                logger.error("添加与会者异常: %s", e)
                return False

        # Step 1: 必选与会者（标准调用，不带 is_optional）
        required_entries = [{"type": "user", "user_id": oid} for oid in attendee_open_ids]
        ok = _post_attendees(required_entries)

        # Step 2: 可选与会者（带 is_optional: True，失败则降级为普通成员）
        opt_ids = [oid for oid in (optional_open_ids or []) if oid not in set(attendee_open_ids)]
        if opt_ids:
            opt_entries_with_flag = [{"type": "user", "user_id": oid, "is_optional": True} for oid in opt_ids]
            if not _post_attendees(opt_entries_with_flag):
                logger.warning("is_optional 标记失败，降级为普通成员再次添加")
                _post_attendees([{"type": "user", "user_id": oid} for oid in opt_ids])

        logger.info("与会者邀请已发送: event=%s, 必选 %d 人, 可选 %d 人", event_id, len(attendee_open_ids), len(opt_ids))
        return ok

    def get_event_attendees(self, calendar_id: str, event_id: str) -> list[dict]:
        """获取日历事件与会者列表及 RSVP 状态（用于轮询拒绝通知）。"""
        try:
            resp = requests.get(
                f"{FEISHU_HOST}/open-apis/calendar/v4/calendars/{calendar_id}/events/{event_id}/attendees",
                headers=self._headers(),
                params={"user_id_type": "open_id"},
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
            if data.get("code") == 0:
                return data.get("data", {}).get("items", [])
        except Exception as e:
            logger.warning("获取与会者状态失败: %s", e)
        return []

    def list_calendar_events(
        self, time_min: datetime, time_max: datetime, title_keyword: str = ""
    ) -> list[dict]:
        """列出应用日历中的事件，可按标题关键词过滤。返回 [{event_id, calendar_id, title, start}]"""
        calendar_id = self._get_or_create_app_calendar_id()
        if not calendar_id:
            return []
        if time_min.tzinfo is None:
            time_min = TZ_SHANGHAI.localize(time_min)
        if time_max.tzinfo is None:
            time_max = TZ_SHANGHAI.localize(time_max)
        try:
            resp = requests.get(
                f"{FEISHU_HOST}/open-apis/calendar/v4/calendars/{calendar_id}/events",
                headers=self._headers(),
                params={
                    "start_time": str(int(time_min.timestamp())),
                    "end_time": str(int(time_max.timestamp())),
                    "page_size": 50,
                },
                timeout=15,
            )
            data = resp.json()
            if data.get("code") != 0:
                logger.warning("查询日历事件失败: %s", data)
                return []
            events = []
            for e in data.get("data", {}).get("items", []):
                title = e.get("summary", "")
                if title_keyword and title_keyword.lower() not in title.lower():
                    continue
                events.append({
                    "event_id": e.get("event_id", ""),
                    "calendar_id": calendar_id,
                    "title": title,
                    "start": e.get("start_time", {}).get("timestamp", ""),
                })
            return events
        except Exception as e:
            logger.error("list_calendar_events 异常: %s", e)
            return []

    def delete_calendar_event(self, calendar_id: str, event_id: str) -> bool:
        """删除（取消）日历事件。成功返回 True。"""
        try:
            resp = requests.delete(
                f"{FEISHU_HOST}/open-apis/calendar/v4/calendars/{calendar_id}/events/{event_id}",
                headers=self._headers(),
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
            if data.get("code") == 0:
                logger.info("日历事件已删除: %s", event_id)
                return True
            logger.error("删除日历事件失败: %s", data)
            return False
        except Exception as e:
            logger.error("删除日历事件异常: %s", e)
            return False

    def query_freebusy(
        self,
        user_open_ids: list[str],
        time_min: datetime,
        time_max: datetime,
    ) -> dict[str, list[dict]]:
        """查询多个用户的忙碌时段。返回 {open_id: [{"start": ISO, "end": ISO}]}"""
        if time_min.tzinfo is None:
            time_min = TZ_SHANGHAI.localize(time_min)
        if time_max.tzinfo is None:
            time_max = TZ_SHANGHAI.localize(time_max)

        result: dict[str, list[dict]] = {}
        for oid in user_open_ids:
            try:
                resp = requests.post(
                    f"{FEISHU_HOST}/open-apis/calendar/v4/freebusy/query",
                    headers=self._headers(),
                    params={"user_id_type": "open_id"},
                    json={
                        "time_min": time_min.isoformat(),
                        "time_max": time_max.isoformat(),
                        "user_id": oid,
                    },
                    timeout=10,
                )
                resp.raise_for_status()
                data = resp.json()
                if data.get("code") == 0:
                    freebusy_list = data.get("data", {}).get("freebusy_list", [])
                    result[oid] = [
                        {"start": fb.get("start_time", ""), "end": fb.get("end_time", "")}
                        for fb in freebusy_list
                    ]
                else:
                    logger.warning("查询 %s 空闲时间失败: code=%s", oid, data.get("code"))
                    result[oid] = []
            except Exception as e:
                logger.warning("查询 %s 空闲时间异常: %s", oid, e)
                result[oid] = []
        return result

    # ------------------------------------------------------------------ #
    # Wiki
    # ------------------------------------------------------------------ #

    def get_wiki_node_space_id(self, node_token: str) -> str | None:
        """通过飞书 API 直接获取 Wiki 节点所在的 space_id（不依赖缓存）"""
        resp = requests.get(
            f"{FEISHU_HOST}/open-apis/wiki/v2/spaces/get_node",
            headers=self._headers(),
            params={"token": node_token, "obj_type": "wiki"},
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.warning("获取 Wiki 节点信息失败: %s", data)
            return None
        return data.get("data", {}).get("node", {}).get("space_id") or None

    def _get_user_token(self) -> str | None:
        """读取 feishu-sync-cli 存储的用户 access_token（用于需要用户权限的操作）"""
        import os as _os
        token_path = _os.path.join(_os.path.expanduser("~"), ".feishu", "access_token.json")
        try:
            with open(token_path) as f:
                data = json.load(f)
            return data.get("access_token") or None
        except Exception:
            return None

    def _get_user_im_token(self) -> str | None:
        """加载用户 IM access_token，过期时自动用 refresh_token 换新。
        token 由 scripts/authorize_user_im.py 生成，存于 .user_im_token.json。
        """
        try:
            if not os.path.exists(_USER_IM_TOKEN_FILE):
                return None
            with open(_USER_IM_TOKEN_FILE, encoding="utf-8") as f:
                data = json.load(f)
            access_token = data.get("access_token")
            if not access_token:
                return None
            # 检查是否在有效期内（留 60s 余量）
            saved_at = data.get("saved_at", 0)
            expires_in = data.get("expires_in", 7200)
            if saved_at and time.time() < saved_at + expires_in - 60:
                return access_token
            # 尝试用 refresh_token 自动换新
            refresh_token = data.get("refresh_token")
            if not refresh_token:
                logger.warning("user_im_token 已过期且无 refresh_token，请重新运行授权脚本")
                return None
            _app_tok_resp = requests.post(
                f"{FEISHU_HOST}/open-apis/auth/v3/app_access_token/internal",
                json={"app_id": self.app_id, "app_secret": self.app_secret},
                timeout=15,
            )
            _app_tok = _app_tok_resp.json().get("app_access_token", "")
            resp = requests.post(
                f"{FEISHU_HOST}/open-apis/authen/v1/oidc/refresh_access_token",
                headers={"Authorization": f"Bearer {_app_tok}", "Content-Type": "application/json"},
                json={"grant_type": "refresh_token", "refresh_token": refresh_token},
                timeout=15,
            )
            result = resp.json()
            if result.get("code") == 0:
                new_data = result.get("data", {})
                new_data["saved_at"] = int(time.time())
                with open(_USER_IM_TOKEN_FILE, "w", encoding="utf-8") as f:
                    json.dump(new_data, f, ensure_ascii=False, indent=2)
                logger.info("user_im_token 已自动刷新")
                return new_data.get("access_token")
            else:
                logger.warning("user_im_token 刷新失败: %s，请重新运行授权脚本", result)
                return None
        except Exception as e:
            logger.warning("_get_user_im_token 失败: %s", e)
            return None

    def _user_im_headers(self) -> dict | None:
        """返回使用用户 IM token 的请求头；无 token 时返回 None。"""
        token = self._get_user_im_token()
        if not token:
            return None
        return {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}

    def get_user_joined_groups(self) -> list[dict]:
        """用用户身份获取用户所在的所有群（含 Bot 不在的群）。返回 [{chat_id, name}]"""
        headers = self._user_im_headers()
        if not headers:
            return []
        groups: list[dict] = []
        page_token = ""
        while True:
            params: dict = {"page_size": 100, "user_id_type": "open_id"}
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/im/v1/chats",
                    headers=headers,
                    params=params,
                    timeout=15,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.warning("get_user_joined_groups 失败: %s", data)
                    break
                for item in data.get("data", {}).get("items", []):
                    if item.get("chat_type") != "p2p":
                        groups.append({"chat_id": item["chat_id"], "name": item.get("name", "")})
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("get_user_joined_groups 异常: %s", e)
                break
        return groups

    def get_group_messages_as_user(
        self, chat_id: str, start_ts: int, end_ts: int, max_msgs: int = 200
    ) -> list[dict]:
        """用用户身份读取群消息（适用于 Bot 不在该群的场景）。"""
        headers = self._user_im_headers()
        if not headers:
            logger.warning("get_group_messages_as_user: 无用户 IM token，请运行授权脚本")
            return []
        messages: list[dict] = []
        page_token = ""
        while len(messages) < max_msgs:
            params: dict = {
                "container_id_type": "chat",
                "container_id": chat_id,
                "start_time": str(start_ts),
                "end_time": str(end_ts),
                "sort_type": "ByCreateTimeAsc",
                "page_size": 50,
            }
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/im/v1/messages",
                    headers=headers,
                    params=params,
                    timeout=15,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.warning(
                        "get_group_messages_as_user [%s] 失败: %s", chat_id, data
                    )
                    break
                for item in data.get("data", {}).get("items", []):
                    try:
                        msg_type = item.get("msg_type", "")
                        if msg_type not in ("text", "post"):
                            continue
                        body = item.get("body", {})
                        content_str = body.get("content", "{}")
                        try:
                            content = json.loads(content_str)
                        except Exception:
                            content = {}
                        if msg_type == "text":
                            text = content.get("text", "").strip()
                        else:
                            text = self._extract_post_text(content)
                        if not text:
                            continue
                        sender_id = item.get("sender", {}).get("id", "")
                        sender_name = self._resolve_sender_name(sender_id)
                        messages.append({
                            "message_id": item.get("message_id", ""),
                            "sender_name": sender_name,
                            "text": text,
                            "ts": int(item.get("create_time", "0")),
                        })
                    except Exception as item_err:
                        logger.warning("跳过单条消息解析失败: %s", item_err)
                        continue
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("get_group_messages_as_user 异常: %s", e)
                break
        return messages

    def move_wiki_page(
        self, space_id: str, node_token: str, target_parent_token: str
    ) -> dict | None:
        """将 Wiki 节点移动到新的父节点下。成功返回新节点信息，失败返回 None。

        移动操作需要用户权限，优先使用 feishu-sync-cli 存储的用户 token。
        正确 API：POST /spaces/{space_id}/nodes/{node_token}/move
        """
        user_token = self._get_user_token()
        if user_token:
            headers = {"Authorization": f"Bearer {user_token}", "Content-Type": "application/json"}
        else:
            headers = self._headers()
            logger.warning("未找到用户 token，使用 tenant token（可能因权限不足失败）")

        resp = requests.post(
            f"{FEISHU_HOST}/open-apis/wiki/v2/spaces/{space_id}/nodes/{node_token}/move",
            headers=headers,
            json={"target_parent_token": target_parent_token},
            timeout=15,
        )
        resp.raise_for_status()
        data = resp.json()
        if data.get("code") != 0:
            logger.error("移动 Wiki 页面失败: %s", data)
            return None
        return data.get("data", {}).get("node", {})

    # ── 文档 @mention 写入 ──────────────────────────────────────────

    # block_type → 文本内容字段名（用于 GET/PATCH docx blocks API）
    _BLOCK_CONTENT_KEYS: dict = {
        2:  "text",
        3:  "heading1",  4:  "heading2",  5:  "heading3",
        6:  "heading4",  7:  "heading5",  8:  "heading6",
        9:  "heading7",  10: "heading8",  11: "heading9",
        12: "bullet",    13: "ordered",
        15: "quote",     17: "todo",
    }

    def _get_docx_blocks(self, document_id: str) -> list[dict]:
        """获取飞书 docx 文档所有 block（自动分页）。"""
        blocks: list[dict] = []
        page_token = ""
        while True:
            params: dict = {"page_size": 500}
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/docx/v1/documents/{document_id}/blocks",
                    headers=self._headers(),
                    params=params,
                    timeout=15,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.error("_get_docx_blocks 失败: %s", data)
                    break
                blocks.extend(data.get("data", {}).get("items", []))
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("_get_docx_blocks 异常: %s", e)
                break
        return blocks

    def apply_mentions_to_wiki_page(
        self,
        node_token: str,
        mention_map: dict,  # {name: open_id}
    ) -> tuple[bool, str]:
        """
        将飞书 wiki 页面中的 @Name 纯文本替换为真正的 mention_user 元素。
        通过 Feishu docx v1 PATCH blocks API 直接修改文档块。
        返回 (成功, 说明文字)。
        """
        if not mention_map:
            return True, "mention_map 为空，无需处理。"

        # 1. 获取 wiki 节点对应的 docx 文档 token
        try:
            resp = requests.get(
                f"{FEISHU_HOST}/open-apis/wiki/v2/spaces/get_node",
                headers=self._headers(),
                params={"token": node_token, "obj_type": "wiki"},
                timeout=10,
            )
            resp.raise_for_status()
            data = resp.json()
        except Exception as e:
            return False, f"获取 wiki 节点信息失败: {e}"

        if data.get("code") != 0:
            return False, f"wiki get_node 失败: code={data.get('code')}, {data.get('msg')}"

        node    = data.get("data", {}).get("node", {})
        doc_id  = node.get("obj_token", "")
        obj_type = node.get("obj_type", "")
        if not doc_id:
            return False, "wiki 节点未返回 obj_token"
        if obj_type != "docx":
            return False, f"文档类型为 {obj_type}，仅支持 docx 格式"

        # 2. 读取所有 blocks
        blocks = self._get_docx_blocks(doc_id)
        if not blocks:
            return False, "未读取到任何 block"

        # 3. 构建 @Name 正则（按名字长度倒序避免子串先匹配）
        names   = sorted(mention_map.keys(), key=len, reverse=True)
        pattern = re.compile(r'@(' + '|'.join(re.escape(n) for n in names) + r')')

        # 4. 优先使用 user token（页面由 feishu-sync-cli 用户 token 创建）
        user_token = self._get_user_token()
        headers = (
            {"Authorization": f"Bearer {user_token}", "Content-Type": "application/json"}
            if user_token else self._headers()
        )

        updated = 0
        errors: list[str] = []

        for block in blocks:
            bt          = block.get("block_type")
            content_key = self._BLOCK_CONTENT_KEYS.get(bt)
            if not content_key:
                continue

            text_content = block.get(content_key, {})
            elements     = text_content.get("elements", [])

            # 检查是否存在含 @Name 的 text_run
            needs_update = any(
                pattern.search(elem.get("text_run", {}).get("content", ""))
                for elem in elements
                if "text_run" in elem
            )
            if not needs_update:
                continue

            # 重新构建 elements，将 @Name 替换为 mention_user
            new_elements: list[dict] = []
            for elem in elements:
                if "text_run" not in elem:
                    new_elements.append(elem)
                    continue
                content = elem["text_run"].get("content", "")
                style   = elem["text_run"].get("text_element_style", {})
                parts   = pattern.split(content)
                for i, part in enumerate(parts):
                    if i % 2 == 0:          # 普通文本段
                        if part:
                            e: dict = {"text_run": {"content": part}}
                            if style:
                                e["text_run"]["text_element_style"] = style
                            new_elements.append(e)
                    else:                   # @Name captured group
                        oid = mention_map.get(part, "")
                        if oid:
                            new_elements.append({"mention_user": {"user_id": oid}})
                        else:
                            new_elements.append({"text_run": {"content": f"@{part}"}})

            # 构造 PATCH body（保留原有 style 等字段，只更新 elements）
            updated_content = {**text_content, "elements": new_elements}
            block_id = block.get("block_id", "")
            try:
                patch_resp = requests.patch(
                    f"{FEISHU_HOST}/open-apis/docx/v1/documents/{doc_id}/blocks/{block_id}",
                    headers=headers,
                    params={"document_revision_id": -1},
                    json={content_key: updated_content},
                    timeout=15,
                )
                patch_data = patch_resp.json()
                if patch_data.get("code") == 0:
                    updated += 1
                    logger.info("apply_mentions: block %s 更新成功", block_id)
                else:
                    msg = f"block {block_id}: code={patch_data.get('code')}, {patch_data.get('msg')}"
                    logger.error("apply_mentions PATCH 失败: %s", msg)
                    errors.append(msg)
            except Exception as e:
                logger.error("apply_mentions PATCH 异常: %s", e)
                errors.append(str(e))

        if updated == 0 and errors:
            return False, f"全部 PATCH 失败：{'; '.join(errors[:3])}"
        return True, f"成功更新 {updated} 个 block" + (f"，{len(errors)} 个失败" if errors else "")

    # ------------------------------------------------------------------ #
    # Bitable
    # ------------------------------------------------------------------ #

    # ── 群消息读取（每日摘要用） ─────────────────────────────────────

    def get_joined_groups(self) -> list[dict]:
        """获取 Bot 所在的所有群组列表。返回 [{chat_id, name}]"""
        groups: list[dict] = []
        page_token = ""
        while True:
            params: dict = {"page_size": 100, "user_id_type": "open_id"}
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/im/v1/chats",
                    headers=self._headers(),
                    params=params,
                    timeout=15,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.warning("get_joined_groups 失败: %s", data)
                    break
                for item in data.get("data", {}).get("items", []):
                    # chat_type 可能为 "group"、None 或缺失，仅排除明确的 p2p 会话
                    if item.get("chat_type") != "p2p":
                        groups.append({"chat_id": item["chat_id"], "name": item.get("name", "")})
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("get_joined_groups 异常: %s", e)
                break
        return groups

    def get_chat_members(self, chat_id: str) -> list[dict]:
        """获取群成员列表。返回 [{open_id, name}]，不依赖消息历史。"""
        members: list[dict] = []
        page_token = ""
        while True:
            params: dict = {"page_size": 100, "member_id_type": "open_id"}
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/im/v1/chats/{chat_id}/members",
                    headers=self._headers(),
                    params=params,
                    timeout=10,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.warning("get_chat_members [%s] 失败: %s", chat_id, data)
                    break
                for item in data.get("data", {}).get("items", []):
                    open_id = item.get("member_id", "")
                    name = item.get("name", "")
                    if open_id:
                        members.append({"open_id": open_id, "name": name})
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("get_chat_members 异常: %s", e)
                break
        return members

    def create_group_chat(self, name: str, member_open_ids: list[str]) -> dict:
        """创建飞书群聊，将 Bot 和指定成员加入。
        返回 {"chat_id": "oc_xxx", "name": "群名", "invite_link": "...", "added": n}，
        失败时 chat_id 为空。
        """
        try:
            # Step 1: 创建群（只建群，不依赖 member_id_list，该参数在 Bot token 下静默失败）
            resp = requests.post(
                f"{FEISHU_HOST}/open-apis/im/v1/chats",
                headers=self._headers(),
                params={"user_id_type": "open_id"},
                json={"name": name},
                timeout=15,
            )
            data = resp.json()
            if data.get("code") != 0:
                logger.error("create_group_chat 失败: %s", data)
                return {"chat_id": "", "name": name, "error": data.get("msg", "unknown")}
            chat_id = data.get("data", {}).get("chat_id", "")

            # Step 2: 批量添加成员（独立接口，Bot token 下有效）
            added = 0
            if member_open_ids and chat_id:
                add_resp = requests.post(
                    f"{FEISHU_HOST}/open-apis/im/v1/chats/{chat_id}/members",
                    headers=self._headers(),
                    params={"member_id_type": "open_id"},
                    json={"id_list": member_open_ids},
                    timeout=15,
                )
                add_data = add_resp.json()
                if add_data.get("code") == 0:
                    added = len(member_open_ids)
                    logger.info("create_group_chat 成员添加成功: %d 人", added)
                else:
                    logger.warning("create_group_chat 添加成员失败: %s", add_data)

            # Step 3: 获取邀请链接
            invite_link = ""
            try:
                link_resp = requests.post(
                    f"{FEISHU_HOST}/open-apis/im/v1/chats/{chat_id}/link",
                    headers=self._headers(),
                    json={},
                    timeout=10,
                )
                link_data = link_resp.json()
                if link_data.get("code") == 0:
                    invite_link = link_data.get("data", {}).get("share_link", "")
            except Exception:
                pass
            return {"chat_id": chat_id, "name": name, "invite_link": invite_link, "added": added}
        except Exception as e:
            logger.error("create_group_chat 异常: %s", e)
            return {"chat_id": "", "name": name, "error": str(e)}

    def delete_group_chat(self, chat_id: str) -> dict:
        """解散飞书群聊（Bot 必须是群主/管理员）。
        返回 {"ok": True} 或 {"ok": False, "error": "..."}。
        """
        try:
            resp = requests.delete(
                f"{FEISHU_HOST}/open-apis/im/v1/chats/{chat_id}",
                headers=self._headers(),
                timeout=15,
            )
            data = resp.json()
            if data.get("code") == 0:
                return {"ok": True}
            return {"ok": False, "error": data.get("msg", "unknown")}
        except Exception as e:
            logger.error("delete_group_chat 异常: %s", e)
            return {"ok": False, "error": str(e)}

    def get_group_messages(
        self, chat_id: str, start_ts: int, end_ts: int, max_msgs: int = 200
    ) -> list[dict]:
        """获取群聊指定时间段内的消息。返回 [{sender_name, text, ts}]"""
        messages: list[dict] = []
        page_token = ""
        while len(messages) < max_msgs:
            params: dict = {
                "container_id_type": "chat",
                "container_id": chat_id,
                "start_time": str(start_ts),
                "end_time": str(end_ts),
                "sort_type": "ByCreateTimeAsc",
                "page_size": 50,
            }
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/im/v1/messages",
                    headers=self._headers(),
                    params=params,
                    timeout=15,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.warning("get_group_messages [%s] 失败: %s", chat_id, data)
                    # 权限错误时尝试用用户 token 重试
                    if not messages:
                        user_msgs = self.get_group_messages_as_user(
                            chat_id, start_ts, end_ts, max_msgs
                        )
                        if user_msgs:
                            return user_msgs
                    break
                for item in data.get("data", {}).get("items", []):
                    try:
                        msg_type = item.get("msg_type", "")
                        if msg_type not in ("text", "post"):
                            continue
                        body = item.get("body", {})
                        content_str = body.get("content", "{}")
                        try:
                            content = json.loads(content_str)
                        except Exception:
                            content = {}
                        if msg_type == "text":
                            text = content.get("text", "").strip()
                        else:
                            text = self._extract_post_text(content)
                        if not text:
                            continue
                        sender_id = item.get("sender", {}).get("id", "")
                        sender_name = self._resolve_sender_name(sender_id)
                        messages.append({
                            "message_id": item.get("message_id", ""),
                            "sender_name": sender_name,
                            "text": text,
                            "ts": int(item.get("create_time", "0")),
                        })
                    except Exception as item_err:
                        logger.warning("跳过单条消息解析失败: %s", item_err)
                        continue
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("get_group_messages 异常: %s", e)
                break
        return messages

    def get_merge_forward_messages(
        self, create_message_id: str, max_msgs: int = 300
    ) -> list[dict]:
        """读取合并转发消息包内的原始消息列表。
        使用 container_id_type=merge_forward_chat，container_id=create_message_id。
        返回 [{sender_name, text, ts, message_id}]
        """
        messages: list[dict] = []
        page_token = ""
        while len(messages) < max_msgs:
            params: dict = {
                "container_id_type": "merge_forward_chat",
                "container_id": create_message_id,
                "sort_type": "ByCreateTimeAsc",
                "page_size": 50,
            }
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/im/v1/messages",
                    headers=self._headers(),
                    params=params,
                    timeout=15,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.warning(
                        "get_merge_forward_messages [%s] 失败: %s",
                        create_message_id, data,
                    )
                    break
                for item in data.get("data", {}).get("items", []):
                    try:
                        msg_type = item.get("msg_type", "")
                        if msg_type not in ("text", "post"):
                            continue
                        body = item.get("body", {})
                        content_str = body.get("content", "{}")
                        try:
                            content = json.loads(content_str)
                        except Exception:
                            content = {}
                        if msg_type == "text":
                            text = content.get("text", "").strip()
                        else:
                            text = self._extract_post_text(content)
                        if not text:
                            continue
                        sender_id = item.get("sender", {}).get("id", "")
                        sender_name = self._resolve_sender_name(sender_id)
                        messages.append({
                            "message_id": item.get("message_id", ""),
                            "sender_name": sender_name,
                            "text": text,
                            "ts": int(item.get("create_time", "0")),
                        })
                    except Exception as item_err:
                        logger.warning("跳过单条消息解析失败: %s", item_err)
                        continue
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("get_merge_forward_messages 异常: %s", e)
                break
        return messages

    def _extract_post_text(self, content: dict) -> str:
        """从 post 消息内容中提取纯文本"""
        for lang in ("zh_cn", "en_us"):
            lang_body = content.get(lang, {})
            if not lang_body:
                continue
            parts: list[str] = []
            title = lang_body.get("title", "")
            if title:
                parts.append(title)
            for paragraph in lang_body.get("content", []):
                for elem in paragraph:
                    tag = elem.get("tag", "")
                    if tag == "text":
                        parts.append(elem.get("text", ""))
                    elif tag == "a":
                        t = elem.get("text", "")
                        h = elem.get("href", "")
                        parts.append(f"{t}({h})" if h and t != h else t or h)
                    elif tag == "at":
                        parts.append(elem.get("user_name", ""))
            return " ".join(p for p in parts if p)
        return ""

    def _resolve_sender_name(self, open_id: str) -> str:
        """把 open_id 转为姓名，先查缓存，再查 API"""
        if not open_id:
            return "未知"
        # Bot/App 的 sender_id 以 cli_ 开头，无法用 /contact/v3/users 查询
        if open_id.startswith("cli_"):
            return "机器人"
        # 先查用户缓存
        for u in self._user_cache:
            if u.get("open_id") == open_id:
                return u.get("name", open_id)
        # 再查 API，捕获异常避免单条消息的名称解析失败中断整个循环
        try:
            info = self.get_user_by_open_id(open_id)
            if info:
                return info.get("name", open_id)
        except Exception:
            pass
        return open_id

    def list_wiki_spaces(self) -> list[dict]:
        """列出所有可访问的飞书 Wiki 空间。返回 [{space_id, name, space_type}]"""
        spaces: list[dict] = []
        page_token = ""
        while True:
            params: dict = {"page_size": 50, "lang": "zh"}
            if page_token:
                params["page_token"] = page_token
            try:
                resp = requests.get(
                    f"{FEISHU_HOST}/open-apis/wiki/v2/spaces",
                    headers=self._headers(),
                    params=params,
                    timeout=15,
                )
                data = resp.json()
                if data.get("code") != 0:
                    logger.warning("list_wiki_spaces 失败: %s", data)
                    break
                for item in data.get("data", {}).get("items", []):
                    spaces.append({
                        "space_id": item.get("space_id", ""),
                        "name": item.get("name", ""),
                        "space_type": item.get("space_type", ""),
                    })
                if not data.get("data", {}).get("has_more"):
                    break
                page_token = data["data"].get("page_token", "")
            except Exception as e:
                logger.error("list_wiki_spaces 异常: %s", e)
                break
        return spaces

    def get_space_root_node_token(self, space_id: str) -> str | None:
        """获取指定 Wiki 空间的根节点 token"""
        try:
            resp = requests.get(
                f"{FEISHU_HOST}/open-apis/wiki/v2/spaces/{space_id}/nodes",
                headers=self._headers(),
                params={"page_size": 1},
                timeout=15,
            )
            data = resp.json()
            items = data.get("data", {}).get("items", [])
            if items:
                return items[0].get("node_token")
        except Exception as e:
            logger.error("get_space_root_node_token 异常: %s", e)
        return None

    # ── Bitable ───────────────────────────────────────────────────────

    def get_bitable_records(
        self, app_token: str, table_id: str, filter_formula: str = ""
    ) -> list[dict]:
        """读取多维表格所有记录（自动分页）"""
        url = f"{FEISHU_HOST}/open-apis/bitable/v1/apps/{app_token}/tables/{table_id}/records"
        records: list[dict] = []
        page_token = ""
        while True:
            params: dict = {"page_size": 100}
            if page_token:
                params["page_token"] = page_token
            if filter_formula:
                params["filter"] = filter_formula
            resp = requests.get(url, headers=self._headers(), params=params, timeout=15)
            resp.raise_for_status()
            data = resp.json()
            if data.get("code") != 0:
                logger.error("读取 Bitable 失败: %s", data)
                break
            items = data.get("data", {}).get("items", [])
            records.extend(items)
            if not data.get("data", {}).get("has_more"):
                break
            page_token = data["data"].get("page_token", "")
        return records
